"""
PPTX Intelligence Slide Integration
====================================
Functions to add intelligence assessment slides to existing PPTX exports.
Integrates with the existing pptx_export.py module.
"""

import json
from typing import Dict, Optional, List, Any
from datetime import datetime

# Note: These functions are designed to integrate with your existing pptx_export.py
# They follow the same patterns and can be called from your existing export flow


def add_intelligence_slide(
    prs,  # python-pptx Presentation object
    site_data: Dict,
    slide_layout_index: int = 5,  # Typically "Title Only" or "Blank"
) -> None:
    """
    Add an intelligence assessment slide to the presentation.
    
    Args:
        prs: python-pptx Presentation object
        site_data: Site data dict containing diagnosis/triage fields
        slide_layout_index: Index of slide layout to use
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
    
    # Get layout
    slide_layout = prs.slide_layouts[slide_layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Intelligence Assessment"
    
    # Parse diagnosis data
    diagnosis_data = {}
    if site_data.get('diagnosis_json'):
        try:
            diagnosis_data = json.loads(site_data['diagnosis_json']) if isinstance(site_data['diagnosis_json'], str) else site_data['diagnosis_json']
        except (json.JSONDecodeError, TypeError):
            pass
    
    # Calculate positions
    left_margin = Inches(0.5)
    top_start = Inches(1.5)
    content_width = Inches(4.5)
    
    # =========================================================================
    # LEFT COLUMN: Timeline Assessment
    # =========================================================================
    
    # Timeline Assessment Header
    timeline_header = slide.shapes.add_textbox(
        left_margin, top_start, content_width, Inches(0.4)
    )
    tf = timeline_header.text_frame
    p = tf.paragraphs[0]
    p.text = "TIMELINE ASSESSMENT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x33, 0x33, 0x33)
    
    # Timeline content
    claimed = site_data.get('claimed_timeline', 'N/A')
    validated = site_data.get('validated_timeline', diagnosis_data.get('validated_timeline', 'N/A'))
    timeline_risk = site_data.get('timeline_risk', diagnosis_data.get('timeline_risk', 'not_assessed'))
    delta = diagnosis_data.get('timeline_delta_months', 0)
    
    timeline_content = slide.shapes.add_textbox(
        left_margin, top_start + Inches(0.5), content_width, Inches(1.5)
    )
    tf = timeline_content.text_frame
    tf.word_wrap = True
    
    # Risk indicator colors
    risk_colors = {
        'on_track': ('✓', RgbColor(0x2e, 0x7d, 0x32)),  # Green
        'at_risk': ('⚠', RgbColor(0xf5, 0x7c, 0x00)),   # Orange
        'not_credible': ('✗', RgbColor(0xc6, 0x28, 0x28)),  # Red
    }
    risk_icon, risk_color = risk_colors.get(timeline_risk, ('•', RgbColor(0x75, 0x75, 0x75)))
    
    lines = [
        f"• Claimed: {claimed}",
        f"• Validated: {validated}",
        f"• Delta: +{delta} months" if delta > 0 else f"• Delta: {delta} months",
        f"• Risk Level: {risk_icon} {timeline_risk.replace('_', ' ').upper()}",
    ]
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        if 'Risk Level' in line:
            p.font.color.rgb = risk_color
            p.font.bold = True
    
    # =========================================================================
    # RIGHT COLUMN: Recommendation & Risks
    # =========================================================================
    
    right_margin = Inches(5.5)
    
    # Recommendation Header
    rec_header = slide.shapes.add_textbox(
        right_margin, top_start, content_width, Inches(0.4)
    )
    tf = rec_header.text_frame
    p = tf.paragraphs[0]
    p.text = "RECOMMENDATION"
    p.font.size = Pt(14)
    p.font.bold = True
    
    # Recommendation value
    recommendation = site_data.get('diagnosis_recommendation', diagnosis_data.get('recommendation', 'NOT_ASSESSED'))
    rec_colors = {
        'GO': RgbColor(0x2e, 0x7d, 0x32),
        'CONDITIONAL_GO': RgbColor(0xf5, 0x7c, 0x00),
        'NO_GO': RgbColor(0xc6, 0x28, 0x28),
    }
    
    rec_box = slide.shapes.add_textbox(
        right_margin, top_start + Inches(0.5), content_width, Inches(0.5)
    )
    tf = rec_box.text_frame
    p = tf.paragraphs[0]
    p.text = recommendation.replace('_', ' ')
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = rec_colors.get(recommendation, RgbColor(0x75, 0x75, 0x75))
    
    # =========================================================================
    # TOP RISKS Section
    # =========================================================================
    
    risks_header = slide.shapes.add_textbox(
        left_margin, top_start + Inches(2.2), Inches(9), Inches(0.4)
    )
    tf = risks_header.text_frame
    p = tf.paragraphs[0]
    p.text = "TOP RISKS"
    p.font.size = Pt(14)
    p.font.bold = True
    
    # Get risks
    top_risks = diagnosis_data.get('top_risks', [])
    if not top_risks and site_data.get('diagnosis_top_risks'):
        top_risks = [r.strip() for r in site_data['diagnosis_top_risks'].split(',')]
    
    risks_box = slide.shapes.add_textbox(
        left_margin, top_start + Inches(2.6), Inches(9), Inches(1.2)
    )
    tf = risks_box.text_frame
    tf.word_wrap = True
    
    for i, risk in enumerate(top_risks[:4]):  # Max 4 risks
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {risk}"
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(0xc6, 0x28, 0x28)
    
    if not top_risks:
        p = tf.paragraphs[0]
        p.text = "• No significant risks identified"
        p.font.size = Pt(11)
    
    # =========================================================================
    # REQUIRED ACTIONS Section
    # =========================================================================
    
    actions_header = slide.shapes.add_textbox(
        left_margin, top_start + Inches(4.0), Inches(9), Inches(0.4)
    )
    tf = actions_header.text_frame
    p = tf.paragraphs[0]
    p.text = "REQUIRED ACTIONS"
    p.font.size = Pt(14)
    p.font.bold = True
    
    # Get actions
    follow_ups = diagnosis_data.get('follow_up_actions', [])
    if not follow_ups and site_data.get('diagnosis_follow_ups'):
        follow_ups = [a.strip() for a in site_data['diagnosis_follow_ups'].split(',')]
    
    actions_box = slide.shapes.add_textbox(
        left_margin, top_start + Inches(4.4), Inches(9), Inches(1.2)
    )
    tf = actions_box.text_frame
    tf.word_wrap = True
    
    for i, action in enumerate(follow_ups[:4]):  # Max 4 actions
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"☐ {action}"
        p.font.size = Pt(11)
    
    if not follow_ups:
        p = tf.paragraphs[0]
        p.text = "☐ Complete full diagnosis"
        p.font.size = Pt(11)
    
    # =========================================================================
    # Footer with metadata
    # =========================================================================
    
    diagnosis_date = site_data.get('diagnosis_date', diagnosis_data.get('diagnosis_date', ''))
    if diagnosis_date:
        footer = slide.shapes.add_textbox(
            left_margin, Inches(6.8), Inches(9), Inches(0.3)
        )
        tf = footer.text_frame
        p = tf.paragraphs[0]
        # Parse and format date
        try:
            dt = datetime.fromisoformat(diagnosis_date.replace('Z', '+00:00'))
            date_str = dt.strftime('%B %d, %Y')
        except:
            date_str = diagnosis_date[:10] if len(diagnosis_date) >= 10 else diagnosis_date
        p.text = f"Intelligence Assessment as of {date_str}"
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RgbColor(0x75, 0x75, 0x75)


def add_utility_assessment_slide(
    prs,
    site_data: Dict,
    utility_intel: Optional[Dict] = None,
    slide_layout_index: int = 5,
) -> None:
    """
    Add a detailed utility assessment slide.
    
    Args:
        prs: python-pptx Presentation object
        site_data: Site data dict
        utility_intel: Optional utility intelligence dict
        slide_layout_index: Index of slide layout to use
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    
    slide_layout = prs.slide_layouts[slide_layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    if title_shape:
        utility_name = site_data.get('utility', 'Utility')
        title_shape.text = f"{utility_name} Assessment"
    
    # Parse diagnosis for utility assessment
    diagnosis_data = {}
    if site_data.get('diagnosis_json'):
        try:
            diagnosis_data = json.loads(site_data['diagnosis_json']) if isinstance(site_data['diagnosis_json'], str) else site_data['diagnosis_json']
        except (json.JSONDecodeError, TypeError):
            pass
    
    utility_assessment = diagnosis_data.get('utility_assessment', {})
    if utility_intel:
        # Merge with provided intel
        utility_assessment.update({
            k: v for k, v in utility_intel.items() 
            if k not in utility_assessment or not utility_assessment[k]
        })
    
    # Positions
    left_margin = Inches(0.5)
    top_start = Inches(1.5)
    
    # Appetite indicator
    appetite = utility_assessment.get('appetite', utility_intel.get('appetite_rating', 'unknown') if utility_intel else 'unknown')
    appetite_colors = {
        'aggressive': ('AGGRESSIVE', RgbColor(0x2e, 0x7d, 0x32)),
        'moderate': ('MODERATE', RgbColor(0xf5, 0x7c, 0x00)),
        'defensive': ('DEFENSIVE', RgbColor(0xc6, 0x28, 0x28)),
    }
    appetite_text, appetite_color = appetite_colors.get(appetite.lower(), (appetite.upper(), RgbColor(0x75, 0x75, 0x75)))
    
    appetite_box = slide.shapes.add_textbox(
        left_margin, top_start, Inches(3), Inches(0.8)
    )
    tf = appetite_box.text_frame
    p = tf.paragraphs[0]
    p.text = "UTILITY APPETITE"
    p.font.size = Pt(11)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = appetite_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = appetite_color
    
    # Capacity position
    capacity_box = slide.shapes.add_textbox(
        Inches(4), top_start, Inches(5.5), Inches(0.8)
    )
    tf = capacity_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPACITY POSITION"
    p.font.size = Pt(11)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = utility_assessment.get('capacity_position', 'Unknown')
    p.font.size = Pt(14)
    
    # Key insight
    key_insight = utility_assessment.get('key_insight', '')
    if key_insight:
        insight_box = slide.shapes.add_textbox(
            left_margin, top_start + Inches(1.2), Inches(9), Inches(0.8)
        )
        tf = insight_box.text_frame
        p = tf.paragraphs[0]
        p.text = "KEY INSIGHT"
        p.font.size = Pt(11)
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = key_insight
        p.font.size = Pt(12)
        p.font.italic = True
    
    # Timeline section
    timeline_box = slide.shapes.add_textbox(
        left_margin, top_start + Inches(2.2), Inches(4.5), Inches(2)
    )
    tf = timeline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "INTERCONNECTION TIMELINE"
    p.font.size = Pt(11)
    p.font.bold = True
    
    realistic_timeline = utility_assessment.get('realistic_timeline', 'Unknown')
    p = tf.add_paragraph()
    p.text = f"• Realistic Timeline: {realistic_timeline}"
    p.font.size = Pt(11)
    
    queue_status = utility_assessment.get('queue_status', '')
    if queue_status:
        p = tf.add_paragraph()
        p.text = f"• Queue Status: {queue_status}"
        p.font.size = Pt(11)
    
    # Recent activity
    recent = utility_assessment.get('recent_activity', '')
    if recent:
        activity_box = slide.shapes.add_textbox(
            Inches(5), top_start + Inches(2.2), Inches(4.5), Inches(2)
        )
        tf = activity_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "RECENT ACTIVITY"
        p.font.size = Pt(11)
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = recent
        p.font.size = Pt(11)


def add_competitive_landscape_slide(
    prs,
    site_data: Dict,
    market_snapshot: Optional[Dict] = None,
    slide_layout_index: int = 5,
) -> None:
    """
    Add a competitive landscape slide.
    
    Args:
        prs: python-pptx Presentation object
        site_data: Site data dict
        market_snapshot: Optional market snapshot dict
        slide_layout_index: Index of slide layout to use
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    
    slide_layout = prs.slide_layouts[slide_layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Competitive Landscape"
    
    # Parse diagnosis data
    diagnosis_data = {}
    if site_data.get('diagnosis_json'):
        try:
            diagnosis_data = json.loads(site_data['diagnosis_json']) if isinstance(site_data['diagnosis_json'], str) else site_data['diagnosis_json']
        except (json.JSONDecodeError, TypeError):
            pass
    
    competitive = diagnosis_data.get('competitive_context', {})
    
    # Positions
    left_margin = Inches(0.5)
    top_start = Inches(1.5)
    
    # Regional projects metric
    projects_box = slide.shapes.add_textbox(
        left_margin, top_start, Inches(3), Inches(1)
    )
    tf = projects_box.text_frame
    p = tf.paragraphs[0]
    p.text = "REGIONAL PROJECTS"
    p.font.size = Pt(11)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = str(competitive.get('regional_projects', '?'))
    p.font.size = Pt(36)
    p.font.bold = True
    
    # Market saturation
    saturation = competitive.get('market_saturation', 'unknown')
    sat_colors = {
        'low': RgbColor(0x2e, 0x7d, 0x32),
        'moderate': RgbColor(0xf5, 0x7c, 0x00),
        'high': RgbColor(0xc6, 0x28, 0x28),
    }
    
    sat_box = slide.shapes.add_textbox(
        Inches(4), top_start, Inches(3), Inches(1)
    )
    tf = sat_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MARKET SATURATION"
    p.font.size = Pt(11)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = saturation.upper()
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = sat_colors.get(saturation.lower(), RgbColor(0x75, 0x75, 0x75))
    
    # Key competitors
    competitors = competitive.get('key_competitors', [])
    if competitors:
        comp_box = slide.shapes.add_textbox(
            left_margin, top_start + Inches(1.5), Inches(4.5), Inches(2.5)
        )
        tf = comp_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "KEY COMPETITORS"
        p.font.size = Pt(11)
        p.font.bold = True
        
        for comp in competitors[:6]:
            p = tf.add_paragraph()
            p.text = f"• {comp}"
            p.font.size = Pt(11)
    
    # Differentiation required
    diff = competitive.get('differentiation_required', '')
    if diff:
        diff_box = slide.shapes.add_textbox(
            Inches(5), top_start + Inches(1.5), Inches(4.5), Inches(2.5)
        )
        tf = diff_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "DIFFERENTIATION REQUIRED"
        p.font.size = Pt(11)
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = diff
        p.font.size = Pt(11)
    
    # Add market snapshot data if available
    if market_snapshot:
        projects = market_snapshot.get('active_projects', [])
        if projects:
            snapshot_box = slide.shapes.add_textbox(
                left_margin, top_start + Inches(4.2), Inches(9), Inches(1.5)
            )
            tf = snapshot_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "MARKET SNAPSHOT - ACTIVE PROJECTS"
            p.font.size = Pt(11)
            p.font.bold = True
            
            for proj in projects[:4]:
                if isinstance(proj, dict):
                    p = tf.add_paragraph()
                    p.text = f"• {proj.get('name', 'Unknown')} - {proj.get('developer', 'Unknown')} ({proj.get('capacity_mw', '?')} MW, {proj.get('status', 'unknown')})"
                    p.font.size = Pt(10)


# =============================================================================
# INTEGRATION HELPER
# =============================================================================

def add_all_intelligence_slides(
    prs,
    site_data: Dict,
    utility_intel: Optional[Dict] = None,
    market_snapshot: Optional[Dict] = None,
    include_utility_detail: bool = True,
    include_competitive: bool = True,
) -> None:
    """
    Add all intelligence-related slides to a presentation.
    
    This is the main entry point for adding intelligence slides.
    Call this from your existing PPTX export flow.
    
    Args:
        prs: python-pptx Presentation object
        site_data: Site data dict with diagnosis fields
        utility_intel: Optional utility intelligence dict
        market_snapshot: Optional market snapshot dict
        include_utility_detail: Whether to add detailed utility slide
        include_competitive: Whether to add competitive landscape slide
    """
    # Always add the main intelligence assessment slide
    add_intelligence_slide(prs, site_data)
    
    # Optionally add utility detail slide
    if include_utility_detail:
        add_utility_assessment_slide(prs, site_data, utility_intel)
    
    # Optionally add competitive landscape slide
    if include_competitive:
        add_competitive_landscape_slide(prs, site_data, market_snapshot)
