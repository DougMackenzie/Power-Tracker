"""
PowerPoint Export Module v3
===========================
Comprehensive site profile export including all diagnostic report sections.

Slide Order:
0. Title (dark blue background - from template)
1. Site Profile (white - from template)
2. Site Boundary (white - from template)
3. Topography (white - from template)
4. Capacity Trajectory (white background, dark header bar)
5. Infrastructure & Critical Path (white background, dark header bar)
6. Thank You (dark blue - moved to end from template)

Diagnostic Sections Captured:
- Site overview with key metrics
- Capacity trajectory (Interconnection vs Generation)
- Critical path phases with study/contract status
- Infrastructure readiness scores
- Risks and opportunities
- Market analysis context
- Score analysis (optional radar chart)
"""

import os
import io
import json
import tempfile
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass, field
from copy import deepcopy

# For chart generation
try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates
    from matplotlib.ticker import MaxNLocator
    import numpy as np
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False


# =============================================================================
# CONFIGURATION
# =============================================================================

JLL_COLORS = {
    'dark_blue': '#1a2b4a',
    'red': '#e31837',
    'light_gray': '#f5f5f5',
    'medium_gray': '#666666',
    'dark_gray': '#333333',
    'teal': '#2b6777',           # Interconnection line
    'tan': '#c9b89d',            # Generation line
    'white': '#ffffff',
    'green': '#2e7d32',
    'amber': '#f9a825',
    'light_blue': '#4a90d9',
}

TEMPLATE_VERSION = "1.0"

TEMPLATE_SLIDES = {
    'title': 0,
    'site_profile': 1,
    'site_boundary': 2,
    'topography': 3,
    'thank_you': 4,
}


# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class PhaseData:
    """Data for a single development phase."""
    phase_num: int
    target_mw: float
    voltage_kv: int
    target_online: str
    screening_study: str = 'Not Started'
    contract_study: str = 'Not Started'
    letter_of_agreement: str = 'Not Started'
    energy_contract: str = 'Not Started'
    transmission_type: str = ''
    substation_type: str = ''
    distance_to_transmission: float = 0.0

    def to_dict(self) -> Dict:
        return {k: v for k, v in self.__dict__.items()}


@dataclass
class CapacityTrajectory:
    """Capacity trajectory data."""
    years: List[int]
    interconnection_mw: List[float]
    generation_mw: List[float]
    available_mw: List[float]

    @classmethod
    def from_dict(cls, data: Dict) -> 'CapacityTrajectory':
        print(f"[DEBUG] CapacityTrajectory.from_dict called with {len(data)} items")
        years, interconnection, generation, available = [], [], [], []
        for year_str, values in sorted(data.items()):
            try:
                year = int(year_str)
                years.append(year)
                # Handle various key formats (interconnection_mw, capacity, ic_mw)
                ic = values.get('interconnection_mw', values.get('capacity', values.get('ic_mw', 0)))
                # Handle various key formats (generation_mw, ramp, gen_mw)
                gen = values.get('generation_mw', values.get('ramp', values.get('gen_mw', 0)))
                
                interconnection.append(ic)
                generation.append(gen)
                available.append(values.get('available_mw', min(ic, gen) if ic and gen else 0))
            except (ValueError, TypeError):
                continue
        print(f"[DEBUG] Parsed years: {years}")
        print(f"[DEBUG] Parsed gen: {generation}")
        return cls(years=years, interconnection_mw=interconnection,
                   generation_mw=generation, available_mw=available)

    @classmethod
    def generate_default(cls, target_mw: float, phase1_mw: float = None,
                         start_year: int = 2028, years: int = 8) -> 'CapacityTrajectory':
        if phase1_mw is None:
            phase1_mw = target_mw * 0.5
        year_list = list(range(start_year, start_year + years))
        interconnection, generation = [], []
        for i in range(len(year_list)):
            ic = phase1_mw * 0.5 if i == 0 else (phase1_mw if i == 1 else target_mw)
            gen = (phase1_mw * 0.3 if i == 0 else
                   (phase1_mw * 0.8 if i == 1 else
                    (phase1_mw if i == 2 else
                     (target_mw * 0.8 if i == 3 else target_mw))))
            interconnection.append(ic)
            generation.append(gen)
        available = [min(ic, gen) for ic, gen in zip(interconnection, generation)]
        return cls(years=year_list, interconnection_mw=interconnection,
                   generation_mw=generation, available_mw=available)


@dataclass
class ScoreAnalysis:
    """Site scoring analysis."""
    overall_score: float = 0
    power_pathway: float = 0
    site_specific: float = 0
    execution: float = 0
    relationship_capital: float = 0
    financial: float = 0

    # Sub-scores
    land_score: float = 0
    queue_score: float = 0
    water_score: float = 0
    fiber_score: float = 0

    @classmethod
    def from_dict(cls, data: Dict):
        """Create ScoreAnalysis from dictionary."""
        if not data:
            return cls()
        # Filter keys to only those in the dataclass
        valid_keys = {k: v for k, v in data.items() if k in cls.__dataclass_fields__}
        return cls(**valid_keys)


@dataclass
class RiskOpportunity:
    """Risk or opportunity item."""
    category: str  # 'risk' or 'opportunity'
    title: str
    description: str
    impact: str = 'Medium'  # Low, Medium, High
    mitigation: str = ''


@dataclass
class SiteProfileData:
    """
    Structured data for Site Profile slide.
    Maps directly to the 17 rows in the Site Profile template table.
    """
    # Header info
    name: str = ''
    state: str = ''
    coordinates: str = ''
    
    # Overview section (left panel)
    overview: str = ''
    observation: str = 'to be completed'
    outstanding: str = 'to be completed'
    
    # Row 1: Location
    nearest_town: str = ''
    distance_to_town: str = ''
    distance_to_airport: str = ''
    airport_name: str = ''
    
    # Row 2: Ownership & Asking Price
    owner_name: str = ''
    willing_to_sell: str = 'TBD'
    asking_price: str = 'TBD'
    
    # Row 3: Size / Shape / Dimensions
    total_acres: float = 0
    phase1_acres: float = 0
    phase2_acres: float = 0
    expandable_acres: float = 0
    shape_description: str = ''
    
    # Row 4: Zoning
    current_zoning: str = 'TBD'
    permits_proposed_use: str = 'TBD'
    zoning_timeline: str = ''
    
    # Row 5: Timing
    site_condition: str = 'Greenfield'  # Greenfield, Brownfield, Existing
    time_to_close: str = 'TBD'
    phase1_delivery: str = ''
    
    # Row 6: Geotechnical and Topography
    soil_type: str = 'TBD'
    bearing_capacity: str = 'TBD'
    topography: str = ''
    slope: str = ''
    
    # Row 7: Environmental, Ecological, Archeological
    environmental_status: str = 'TBD'
    phase1_esa: str = ''
    ecological_concerns: str = ''
    archeological: str = ''
    
    # Row 8: Wetlands and Jurisdictional Water
    wetlands_present: bool = False
    wetlands_acres: float = 0
    wetlands_avoidable: bool = True
    jurisdictional_water: str = ''
    
    # Row 9: Disaster
    flood_zone: str = 'TBD'
    seismic_risk: str = 'Low'
    hurricane_risk: str = 'Low'
    tornado_risk: str = ''
    
    # Row 10: Easements
    easements: str = ''
    right_of_way: str = 'TBD'
    
    # Row 11: Electricity
    electric_utility: str = ''
    voltage_kv: int = 0
    transmission_line: str = ''
    estimated_capacity_mw: float = 0
    distance_to_transmission: str = ''
    
    # Row 12: Water
    water_service: str = ''
    water_provider: str = ''
    water_capacity_gpd: str = ''
    water_line_size: str = ''
    
    # Row 13: Wastewater
    wastewater_solution: str = ''
    wastewater_capacity_gpd: str = ''
    wastewater_provider: str = ''
    
    # Row 14: Telecom
    fiber_provider: str = ''
    fiber_capacity: str = ''
    fiber_distance: str = ''
    lit_building_distance: str = ''
    
    # Row 15: Gas
    gas_provider: str = ''
    gas_capacity: str = ''
    gas_line_size: str = ''
    
    # Row 16: Transportation
    highway_distance: str = ''
    highway_name: str = ''
    airport_distance: str = ''
    rail_access: str = ''
    
    # Row 17: Labor
    workforce_radius: str = '30-mile'
    workforce_population: str = ''
    unemployment_rate: str = ''
    
    # Ratings (1=Desirable, 2=Acceptable, 3=Marginal, 4=Fatal Flaw, 0=NR)
    ratings: Dict[str, int] = field(default_factory=dict)
    
    # --- Program Tracker Fields (Two-Way Sync) ---
    client: str = ''
    total_fee_potential: float = 0.0
    contract_status: str = 'No'
    site_control_stage: int = 1
    power_stage: int = 1
    marketing_stage: int = 1
    buyer_stage: int = 1
    zoning_stage: int = 1
    water_stage: int = 1
    incentives_stage: int = 1
    probability: float = 0.0
    weighted_fee: float = 0.0
    tracker_notes: str = ''
    
    def to_description_dict(self) -> Dict[str, str]:
        """
        Convert to dictionary of description texts for each row.
        Keys match the Item column in the template.
        """
        descriptions = {}
        
        # Row 1: Location
        loc_parts = []
        if self.nearest_town:
            loc_parts.append(f"Nearest Town: {self.nearest_town}")
            if self.distance_to_town:
                loc_parts[-1] += f" – {self.distance_to_town}"
        if self.airport_name and self.distance_to_airport:
            loc_parts.append(f"Distance to {self.airport_name}: {self.distance_to_airport}")
        descriptions['Location'] = ". ".join(loc_parts) if loc_parts else "TBD"
        
        # Row 2: Ownership
        own_parts = []
        if self.owner_name:
            own_parts.append(f"Owner(s): {self.owner_name}")
        own_parts.append(f"Confirmed Willingness to Sell: {self.willing_to_sell}")
        if self.asking_price and self.asking_price != 'TBD':
            own_parts.append(f"Asking Price: {self.asking_price}")
        descriptions['Ownership & Asking Price'] = ". ".join(own_parts)
        
        # Row 3: Size
        size_parts = []
        if self.total_acres:
            size_parts.append(f"Total Size: {self.total_acres:,.0f} acres")
        if self.phase1_acres:
            size_parts.append(f"Phase I: {self.phase1_acres:,.0f} acres")
        if self.phase2_acres:
            size_parts.append(f"Phase II: {self.phase2_acres:,.0f} acres")
        if self.expandable_acres:
            size_parts.append(f"expandable to {self.expandable_acres:,.0f} acres")
        if self.shape_description:
            size_parts.append(self.shape_description)
        descriptions['Size / Shape / Dimensions'] = ", ".join(size_parts) if size_parts else "TBD"
        
        # Row 4: Zoning
        zoning_parts = [f"Zoning: {self.current_zoning}"]
        zoning_parts.append(f"Permits Proposed Use: {self.permits_proposed_use}")
        if self.zoning_timeline:
            zoning_parts.append(self.zoning_timeline)
        descriptions['Zoning'] = ". ".join(zoning_parts)
        
        # Row 5: Timing
        timing_parts = [f"Site condition: {self.site_condition}"]
        timing_parts.append(f"Time to Close: {self.time_to_close}")
        if self.phase1_delivery:
            timing_parts.append(f"Phase 1 Delivery: {self.phase1_delivery}")
        descriptions['Timing'] = ". ".join(timing_parts)
        
        # Row 6: Geotech
        geo_parts = [f"Soil type and bearing capacity: {self.soil_type}"]
        if self.topography:
            geo_parts.append(f"Topography: {self.topography}")
        if self.slope:
            geo_parts.append(f"Slope: {self.slope}")
        descriptions['Geotechnical and Topography'] = ". ".join(geo_parts)
        
        # Row 7: Environmental
        env_parts = [f"Environmental: {self.environmental_status}"]
        if self.phase1_esa:
            env_parts.append(f"Phase I ESA: {self.phase1_esa}")
        if self.ecological_concerns:
            env_parts.append(f"Ecological: {self.ecological_concerns}")
        if self.archeological:
            env_parts.append(f"Archeological: {self.archeological}")
        descriptions['Environmental, Ecological, Archeological'] = ". ".join(env_parts)
        
        # Row 8: Wetlands
        if self.wetlands_present:
            wetland_text = f"Wetlands: Present on site ({self.wetlands_acres:.1f} acres)"
            if self.wetlands_avoidable:
                wetland_text += " but able to avoid for development"
        else:
            wetland_text = "Wetlands: None identified"
        if self.jurisdictional_water:
            wetland_text += f". {self.jurisdictional_water}"
        descriptions['Wetlands and Jurisdictional Water'] = wetland_text
        
        # Row 9: Disaster
        disaster_parts = [f"Flood: {self.flood_zone.strip() if self.flood_zone else ''}"]
        disaster_parts.append(f"Seismic: {self.seismic_risk.strip() if self.seismic_risk else ''}")
        disaster_parts.append(f"Hurricane/Weather: {self.hurricane_risk.strip() if self.hurricane_risk else ''}")
        if self.tornado_risk:
            disaster_parts.append(f"Tornado: {self.tornado_risk.strip()}")
        descriptions['Disaster'] = ". ".join(disaster_parts)
        
        # Row 10: Easements
        ease_parts = []
        if self.easements:
            ease_parts.append(f"Easements: {self.easements}")
        ease_parts.append(f"Right of Way: {self.right_of_way}")
        descriptions['Easements'] = ". ".join(ease_parts) if ease_parts else "TBD"
        
        # Row 11: Electricity
        elec_parts = []
        if self.electric_utility and self.transmission_line:
            elec_parts.append(f"Electric Service to Site: {self.electric_utility} {self.transmission_line}")
        elif self.electric_utility:
            elec_parts.append(f"Electric Service to Site: {self.electric_utility}")
        if self.estimated_capacity_mw:
            mw = self.estimated_capacity_mw
            cap_str = f"{mw/1000:.1f}GW+" if mw >= 1000 else f"{mw:.0f}MW"
            elec_parts.append(f"Estimated Capacity: {cap_str}")
        if self.distance_to_transmission:
            elec_parts.append(f"Distance to transmission: {self.distance_to_transmission}")
        descriptions['Electricity'] = ". ".join(elec_parts) if elec_parts else "TBD"
        
        # Row 12: Water
        water_parts = []
        if self.water_service:
            water_parts.append(f"Current Service to Site: {self.water_service}")
        if self.water_provider:
            water_parts.append(f"Provider: {self.water_provider}")
        if self.water_capacity_gpd:
            water_parts.append(f"Capacity: {self.water_capacity_gpd}")
        descriptions['Water'] = ". ".join(water_parts) if water_parts else "TBD"
        
        # Row 13: Wastewater
        ww_parts = []
        if self.wastewater_solution:
            ww_parts.append(f"Discharge Solution: {self.wastewater_solution}")
        if self.wastewater_capacity_gpd:
            ww_parts.append(f"Capacity: {self.wastewater_capacity_gpd}")
        if self.wastewater_provider:
            ww_parts.append(f"Provider: {self.wastewater_provider}")
        descriptions['Wastewater'] = ". ".join(ww_parts) if ww_parts else "TBD"
        
        # Row 14: Telecom
        fiber_parts = []
        if self.fiber_provider:
            fiber_parts.append(f"Current Service to Site: {self.fiber_provider}")
        if self.fiber_capacity:
            fiber_parts.append(f"Capacity: {self.fiber_capacity}")
        if self.lit_building_distance:
            fiber_parts.append(f"Distance to lit building: {self.lit_building_distance}")
        descriptions['Telecom'] = ". ".join(fiber_parts) if fiber_parts else "TBD"
        
        # Row 15: Gas
        gas_parts = []
        if self.gas_provider:
            gas_parts.append(f"Provider: {self.gas_provider}")
        if self.gas_capacity:
            gas_parts.append(f"Capacity: {self.gas_capacity}")
        if self.gas_line_size:
            gas_parts.append(f"Line size: {self.gas_line_size}")
        descriptions['Gas'] = ". ".join(gas_parts) if gas_parts else "TBD"
        
        # Row 16: Transportation
        trans_parts = []
        if self.highway_name and self.highway_distance:
            trans_parts.append(f"Hwy: {self.highway_distance} to {self.highway_name}")
        if self.airport_distance:
            trans_parts.append(f"Airport: {self.airport_distance}")
        if self.rail_access:
            trans_parts.append(f"Rail: {self.rail_access}")
        descriptions['Transportation'] = ". ".join(trans_parts) if trans_parts else "TBD"
        
        # Row 17: Labor
        labor_parts = []
        if self.workforce_population:
            labor_parts.append(f"Workforce Within {self.workforce_radius} Radius: {self.workforce_population}")
        if self.unemployment_rate:
            labor_parts.append(f"Unemployment: {self.unemployment_rate}")
        descriptions['Labor'] = ". ".join(labor_parts) if labor_parts else "TBD"
        
        return descriptions
    
    @classmethod
    def from_dict(cls, data: Dict) -> 'SiteProfileData':
        """Create from a dictionary (e.g., extracted from website)."""
        profile = cls()
        for key, value in data.items():
            if hasattr(profile, key):
                setattr(profile, key, value)
        return profile
    """
    Market analysis data based on state analysis framework.
    
    Note: Queue time refers to GENERATION interconnection (bringing new power plants online),
    NOT large load interconnection for data center customers.
    """
    
    # Site's state info
    state_code: str = ''
    state_name: str = ''
    
    # ISO/Utility Profile
    primary_iso: str = ''
    regulatory_structure: str = ''  # regulated, deregulated, hybrid
    utility_name: str = ''
    avg_queue_time_months: int = 36  # Generation interconnection queue, NOT large load
    avg_industrial_rate: float = 0.0  # $/kWh
    renewable_percentage: float = 0
    
    # State scores (from StateProfile)
    overall_score: int = 0
    regulatory_score: int = 0
    transmission_score: int = 0
    power_score: int = 0
    water_score: int = 0
    business_score: int = 0
    ecosystem_score: int = 0
    tier: int = 3
    
    # Competitive landscape
    existing_dc_mw: int = 0
    hyperscaler_presence: List[str] = field(default_factory=list)
    fiber_density: str = ''
    
    # State SWOT
    strengths: List[str] = field(default_factory=list)
    weaknesses: List[str] = field(default_factory=list)
    opportunities: List[str] = field(default_factory=list)
    threats: List[str] = field(default_factory=list)
    
    # Comparison states (state_code: {name, queue_months, rate, score})
    comparison_states: Dict[str, Dict] = field(default_factory=dict)
    
    # Data center incentives
    incentives: List[str] = field(default_factory=list)
    
    # Known constraints
    known_moratoria: List[str] = field(default_factory=list)
    
    @classmethod
    def from_state_profile(cls, profile, utility_name: str = '', 
                           comparison_profiles: List = None) -> 'MarketAnalysis':
        """Create MarketAnalysis from StateProfile."""
        comparisons = {}
        if comparison_profiles:
            for cp in comparison_profiles:
                comparisons[cp.state_code] = {
                    'name': cp.state_name,
                    'queue_months': cp.avg_queue_time_months,
                    'rate': cp.avg_industrial_rate,
                    'score': cp.overall_score,
                    'tier': cp.tier,
                }
        
        return cls(
            state_code=profile.state_code,
            state_name=profile.state_name,
            primary_iso=profile.primary_iso,
            regulatory_structure=profile.regulatory_structure,
            utility_name=utility_name or profile.utility_type,
            avg_queue_time_months=profile.avg_queue_time_months,
            avg_industrial_rate=profile.avg_industrial_rate,
            renewable_percentage=profile.renewable_percentage,
            overall_score=profile.overall_score,
            regulatory_score=profile.regulatory_score,
            transmission_score=profile.transmission_score,
            power_score=profile.power_score,
            water_score=profile.water_score,
            business_score=profile.business_score,
            ecosystem_score=profile.ecosystem_score,
            tier=profile.tier,
            existing_dc_mw=profile.existing_dc_mw,
            hyperscaler_presence=profile.hyperscaler_presence,
            fiber_density=profile.fiber_density,
            strengths=profile.strengths,
            weaknesses=profile.weaknesses,
            opportunities=profile.opportunities,
            threats=profile.threats,
            comparison_states=comparisons,
            incentives=profile.data_center_incentives,
            known_moratoria=profile.known_moratoria,
        )


@dataclass
class ExportConfig:
    """Configuration for PPTX export."""
    include_capacity_trajectory: bool = True
    include_infrastructure: bool = True
    include_site_boundary: bool = True
    include_topography: bool = True
    include_score_analysis: bool = True
    include_market_analysis: bool = True
    contact_name: str = ""
    contact_title: str = ""
    contact_phone: str = ""
    contact_email: str = ""
    custom_images: Dict[str, str] = field(default_factory=dict)
    chart_subtitle: str = "Anticipated Power Ramp (Pending Contract Study Results)"


# =============================================================================
# CHART GENERATION
# =============================================================================

def generate_capacity_trajectory_chart(
    trajectory: CapacityTrajectory,
    site_name: str,
    output_path: str,
    phases: List[PhaseData] = None,
    width: float = 11,
    height: float = 5.5,
    subtitle: str = "Anticipated Power Ramp (Pending Contract Study Results)",
) -> str:
    """Generate capacity trajectory chart with Interconnection vs Generation lines."""
    if not MATPLOTLIB_AVAILABLE:
        raise ImportError("matplotlib required")

    fig, ax = plt.subplots(figsize=(width, height), facecolor='white')
    ax.set_facecolor('white')

    dates = [datetime(year, 1, 1) for year in trajectory.years]

    # Interconnection (teal)
    ax.plot(dates, trajectory.interconnection_mw, color=JLL_COLORS['teal'],
            linewidth=2.5, label='Utility Interconnection Rating', marker='o', markersize=4, zorder=3)

    # Generation (tan)
    ax.plot(dates, trajectory.generation_mw, color=JLL_COLORS['tan'],
            linewidth=2.5, label='Utility Generation Capacity', marker='s', markersize=4, zorder=2)

    # Max capacity annotation
    max_ic = max(trajectory.interconnection_mw)
    max_idx = trajectory.interconnection_mw.index(max_ic)
    label = f'Full {max_ic/1000:.1f}GW' if max_ic >= 1000 else f'Full {max_ic:.0f}MW'
    ax.annotate(label, xy=(dates[max_idx], max_ic), xytext=(dates[max_idx], max_ic * 1.05),
                fontsize=11, fontweight='bold', ha='center', color=JLL_COLORS['dark_blue'])

    # Styling
    ax.set_ylabel('Capacity (MW)', fontsize=12, fontweight='bold', color=JLL_COLORS['dark_blue'])
    ax.set_title('Capacity Trajectory', fontsize=16, fontweight='bold',
                 color=JLL_COLORS['dark_blue'], pad=20, loc='left')


    ax.xaxis.set_major_locator(mdates.YearLocator())
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%Y'))
    max_y = max(max(trajectory.interconnection_mw), max(trajectory.generation_mw)) * 1.15
    ax.set_ylim(0, max_y)
    ax.yaxis.set_major_locator(MaxNLocator(integer=True, nbins=6))

    ax.grid(True, axis='y', linestyle='-', alpha=0.3, color=JLL_COLORS['medium_gray'])
    ax.grid(True, axis='x', linestyle='-', alpha=0.2, color=JLL_COLORS['medium_gray'])
    ax.legend(loc='upper left', frameon=True, framealpha=0.9, fontsize=10)

    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color(JLL_COLORS['medium_gray'])
    ax.spines['bottom'].set_color(JLL_COLORS['medium_gray'])
    ax.tick_params(colors=JLL_COLORS['dark_gray'])

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close()
    return output_path


def generate_critical_path_chart(
    phases: List[PhaseData],
    site_data: Dict,
    output_path: str,
    width: float = 6,
    height: float = 4,
) -> str:
    """Generate Infrastructure Readiness bar chart only."""
    if not MATPLOTLIB_AVAILABLE:
        raise ImportError("matplotlib required")

    fig, ax = plt.subplots(figsize=(width, height), facecolor='white')
    
    ax.set_title('Infrastructure Readiness', fontsize=14, fontweight='bold',
                  color=JLL_COLORS['dark_blue'], loc='left', pad=10)

    categories = [
        ('Power', site_data.get('power_stage', 1), 4),
        ('Site Control', site_data.get('site_control_stage', 1), 4),
        ('Zoning', site_data.get('zoning_stage', 1), 3),
        ('Water', site_data.get('water_stage', 1), 4),
        ('Fiber', 3 if site_data.get('fiber_available') else 1, 4),
        ('Environmental', 4 if site_data.get('environmental_complete') else 2, 4),
    ]

    y_positions = list(range(len(categories) - 1, -1, -1))
    for i, (name, stage, max_stage) in enumerate(categories):
        pct = (stage / max_stage) * 100
        ax.barh(y_positions[i], 100, height=0.6, color=JLL_COLORS['light_gray'], zorder=1)
        color = JLL_COLORS['green'] if pct >= 75 else (JLL_COLORS['amber'] if pct >= 50 else JLL_COLORS['teal'])
        ax.barh(y_positions[i], pct, height=0.6, color=color, zorder=2)
        ax.text(pct + 2, y_positions[i], f'{pct:.0f}%', va='center', fontsize=10,
                color=JLL_COLORS['dark_blue'], fontweight='bold')

    ax.set_yticks(y_positions)
    ax.set_yticklabels([c[0] for c in categories], fontsize=10)
    ax.set_xlim(0, 115)
    for spine in ['top', 'right', 'bottom']:
        ax.spines[spine].set_visible(False)
    ax.tick_params(bottom=False, labelbottom=False)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close()
    return output_path


def generate_score_radar_chart(scores: ScoreAnalysis, site_name: str, output_path: str,
                                width: float = 6, height: float = 6) -> str:
    """Generate radar chart for site scores only."""
    if not MATPLOTLIB_AVAILABLE:
        raise ImportError("matplotlib required")

    categories = ['Power\nPathway', 'Site\nSpecific', 'Execution', 'Relationship\nCapital', 'Financial']
    values = [scores.power_pathway, scores.site_specific, scores.execution,
              scores.relationship_capital, scores.financial]

    angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
    values_closed = values + values[:1]
    angles_closed = angles + angles[:1]

    fig, ax = plt.subplots(figsize=(width, height), subplot_kw=dict(polar=True), facecolor='white')
    
    # Fill area
    ax.fill(angles_closed, values_closed, color=JLL_COLORS['teal'], alpha=0.25)
    ax.plot(angles_closed, values_closed, color=JLL_COLORS['teal'], linewidth=2.5, marker='o', markersize=6)
    
    # Add value labels
    for angle, value in zip(angles, values):
        ax.annotate(f'{value:.0f}', xy=(angle, value), xytext=(angle, value + 10),
                   ha='center', fontsize=10, fontweight='bold', color=JLL_COLORS['dark_blue'])
    
    ax.set_xticks(angles)
    ax.set_xticklabels(categories, fontsize=10, color=JLL_COLORS['dark_gray'])
    ax.set_ylim(0, 100)
    ax.set_yticks([25, 50, 75, 100])
    ax.set_yticklabels(['25', '50', '75', '100'], fontsize=8, color=JLL_COLORS['medium_gray'])
    ax.grid(True, color=JLL_COLORS['medium_gray'], alpha=0.3)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close()
    return output_path


def generate_score_summary_chart(scores: ScoreAnalysis, site_name: str, output_path: str,
                                  width: float = 11, height: float = 6) -> str:
    """Generate comprehensive score summary with radar + details."""
    if not MATPLOTLIB_AVAILABLE:
        raise ImportError("matplotlib required")

    fig = plt.figure(figsize=(width, height), facecolor='white')
    gs = fig.add_gridspec(1, 2, width_ratios=[1.2, 1], wspace=0.3)

    # Left: Radar chart
    ax1 = fig.add_subplot(gs[0, 0], polar=True)
    
    categories = ['Power\nPathway', 'Site\nSpecific', 'Execution', 'Relationship\nCapital', 'Financial']
    values = [scores.power_pathway, scores.site_specific, scores.execution,
              scores.relationship_capital, scores.financial]

    angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
    values_closed = values + values[:1]
    angles_closed = angles + angles[:1]

    ax1.fill(angles_closed, values_closed, color=JLL_COLORS['teal'], alpha=0.25)
    ax1.plot(angles_closed, values_closed, color=JLL_COLORS['teal'], linewidth=2.5, marker='o', markersize=6)
    
    for angle, value in zip(angles, values):
        ax1.annotate(f'{value:.0f}', xy=(angle, value), xytext=(angle, value + 10),
                    ha='center', fontsize=10, fontweight='bold', color=JLL_COLORS['dark_blue'])

    ax1.set_xticks(angles)
    ax1.set_xticklabels(categories, fontsize=10, color=JLL_COLORS['dark_gray'])
    ax1.set_ylim(0, 100)
    ax1.set_yticks([25, 50, 75, 100])
    ax1.grid(True, color=JLL_COLORS['medium_gray'], alpha=0.3)

    # Right: Score breakdown
    ax2 = fig.add_subplot(gs[0, 1])
    ax2.set_facecolor('white')
    ax2.axis('off')

    # Overall score display
    ax2.text(0.5, 0.95, 'Overall Score', fontsize=14, fontweight='bold',
            color=JLL_COLORS['dark_blue'], ha='center', transform=ax2.transAxes)
    
    # Score circle
    score_color = JLL_COLORS['green'] if scores.overall_score >= 70 else (
        JLL_COLORS['amber'] if scores.overall_score >= 50 else JLL_COLORS['red'])
    circle = plt.Circle((0.5, 0.75), 0.15, color=score_color, alpha=0.2, transform=ax2.transAxes)
    ax2.add_patch(circle)
    ax2.text(0.5, 0.75, f'{scores.overall_score:.0f}', fontsize=36, fontweight='bold',
            color=score_color, ha='center', va='center', transform=ax2.transAxes)

    # Sub-scores
    sub_scores = [
        ('Power Pathway', scores.power_pathway, 0.30),
        ('Site Specific', scores.site_specific, 0.10),
        ('Execution', scores.execution, 0.20),
        ('Relationships', scores.relationship_capital, 0.35),
        ('Financial', scores.financial, 0.05),
    ]

    y_pos = 0.50
    ax2.text(0.1, y_pos + 0.05, 'Category', fontsize=10, fontweight='bold',
            color=JLL_COLORS['dark_blue'], transform=ax2.transAxes)
    ax2.text(0.55, y_pos + 0.05, 'Score', fontsize=10, fontweight='bold',
            color=JLL_COLORS['dark_blue'], ha='center', transform=ax2.transAxes)
    ax2.text(0.8, y_pos + 0.05, 'Weight', fontsize=10, fontweight='bold',
            color=JLL_COLORS['dark_blue'], ha='center', transform=ax2.transAxes)

    for name, score, weight in sub_scores:
        y_pos -= 0.08
        ax2.text(0.1, y_pos, name, fontsize=10, color=JLL_COLORS['dark_gray'], transform=ax2.transAxes)
        
        # Mini bar
        bar_width = score / 100 * 0.25
        ax2.barh(y_pos, bar_width, height=0.04, left=0.4, color=JLL_COLORS['teal'],
                transform=ax2.transAxes, zorder=2)
        ax2.barh(y_pos, 0.25, height=0.04, left=0.4, color=JLL_COLORS['light_gray'],
                transform=ax2.transAxes, zorder=1)
        
        ax2.text(0.68, y_pos, f'{score:.0f}', fontsize=10, va='center',
                color=JLL_COLORS['dark_blue'], transform=ax2.transAxes)
        ax2.text(0.8, y_pos, f'{weight*100:.0f}%', fontsize=10, va='center', ha='center',
                color=JLL_COLORS['medium_gray'], transform=ax2.transAxes)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close()
    return output_path


def generate_market_analysis_chart(market_data: Dict, site_name: str, output_path: str,
                                    width: float = 6, height: float = 4) -> str:
    """
    Generate market analysis chart (State Comparison only).
    """
    if not MATPLOTLIB_AVAILABLE:
        raise ImportError("matplotlib required")

    fig, ax1 = plt.subplots(figsize=(width, height), facecolor='white')
    
    ax1.set_title('New Generation Capacity: State Comparison', fontsize=11, fontweight='bold',
                  color=JLL_COLORS['dark_blue'], loc='left', pad=8)
    
    # Get state data
    state_code = market_data.get('state_code', 'OK')
    # state_name = market_data.get('state_name', 'State') # Not used in this simplified chart
    queue_months = market_data.get('avg_queue_time_months', 36)
    rate = market_data.get('avg_industrial_rate', 0.06)
    comparisons = market_data.get('comparison_states', {})
    
    # Build comparison data
    states = [state_code]
    queues = [queue_months]
    rates = [rate * 100]  # Convert to cents/kWh
    colors = [JLL_COLORS['teal']] # This will be the color for the current state's queue bar
    
    for code, data in comparisons.items():
        states.append(code)
        queues.append(data.get('queue_months', 36))
        rates.append(data.get('rate', 0.06) * 100)
        colors.append(JLL_COLORS['light_gray']) # Comparison states queue bars
    
    # Create dual bar chart
    x = np.arange(len(states))
    bar_width = 0.35
    
    # Queue bars: current state is teal, others are light_gray
    bars1 = ax1.bar(x - bar_width/2, queues, bar_width, label='Gen. Queue (months)', 
                    color=[JLL_COLORS['teal'] if i == 0 else JLL_COLORS['light_gray'] for i in range(len(states))], alpha=0.8)
    
    ax1_twin = ax1.twinx()
    # Rate bars: all are tan
    bars2 = ax1_twin.bar(x + bar_width/2, rates, bar_width, label='Power Cost (¢/kWh)',
                         color=[JLL_COLORS['tan'] for _ in range(len(states))], alpha=0.7)
    
    ax1.set_xticks(x)
    ax1.set_xticklabels(states, fontsize=10)
    ax1.set_ylabel('Gen. Interconnection (months)', fontsize=9, color=JLL_COLORS['teal'])
    ax1_twin.set_ylabel('Power Cost (¢/kWh)', fontsize=9, color=JLL_COLORS['tan'])
    
    # Add value labels
    for bar, val in zip(bars1, queues):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1, f'{val}',
                ha='center', fontsize=8, color=JLL_COLORS['dark_blue'])
    for bar, val in zip(bars2, rates):
        ax1_twin.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.1, f'{val:.1f}¢',
                     ha='center', fontsize=8, color=JLL_COLORS['dark_gray'])
    
    ax1.spines['top'].set_visible(False)
    ax1_twin.spines['top'].set_visible(False)
    
    # Combine legends and place at bottom
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax1_twin.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper center', 
               bbox_to_anchor=(0.5, -0.15), ncol=2, fontsize=8, frameon=False)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close()
    return output_path


# =============================================================================
# SLIDE HELPERS
# =============================================================================

def set_slide_background_white(slide, Inches, RGBColor):
    """Set slide background to white."""
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    
    # Access the slide's spTree and add a background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def add_header_bar(slide, title_text: str, Inches, Pt, RGBColor):
    """Add dark blue header bar."""
    from pptx.enum.shapes import MSO_SHAPE

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     Inches(13.33), Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor.from_string(JLL_COLORS['dark_blue'][1:])
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(10), Inches(0.4))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string('ffffff')


def add_footer(slide, page_num: int, Inches, Pt, RGBColor):
    """Add JLL footer."""
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.2), Inches(10), Inches(0.25))
    p = footer.text_frame.paragraphs[0]
    p.text = f"© {datetime.now().year} Jones Lang LaSalle IP, Inc. All rights reserved  |  {page_num}"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor.from_string(JLL_COLORS['medium_gray'][1:])


def find_and_replace_text(shape, replacements: Dict[str, str]) -> bool:
    """Find and replace text in shape."""
    if not shape.has_text_frame:
        return False
    replaced = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    replaced = True
    return replaced


def update_overview_textbox(shape, site_data: Dict, profile: Optional['SiteProfileData'] = None):
    """
    Update the Overview/Observation/Outstanding text box.
    This text box has a specific structure with labels and content in separate runs.
    Also fixes any red coloring to use theme color.
    """
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.dml.color import RGBColor
    
    if not shape.has_text_frame:
        return False
    
    updated = False
    
    for para in shape.text_frame.paragraphs:
        para_text = para.text.lower()
        
        # Find which field this paragraph is for
        if 'overview' in para_text:
            new_content = None
            if profile and profile.overview:
                new_content = profile.overview
            elif site_data.get('overview'):
                new_content = site_data['overview']
            
            if new_content and len(para.runs) >= 2:
                # Run 1 has the content (Run 0 has "Overview:  ")
                para.runs[1].text = new_content
                # Fix color - use tan/light brown
                para.runs[1].font.color.rgb = RGBColor.from_string(JLL_COLORS['tan'][1:])
                updated = True
        
        elif 'observation' in para_text:
            new_content = None
            if profile and profile.observation and profile.observation != 'to be completed':
                new_content = profile.observation
            elif site_data.get('observation') and site_data['observation'] != 'to be completed':
                new_content = site_data['observation']
            
            if new_content and len(para.runs) >= 2:
                para.runs[1].text = new_content
                # Use tan/light brown
                para.runs[1].font.color.rgb = RGBColor.from_string(JLL_COLORS['tan'][1:])
                updated = True
        
        elif 'outstanding' in para_text:
            new_content = None
            if profile and profile.outstanding and profile.outstanding != 'to be completed':
                new_content = profile.outstanding
            elif site_data.get('outstanding') and site_data['outstanding'] != 'to be completed':
                new_content = site_data['outstanding']
            
            if new_content and len(para.runs) >= 2:
                para.runs[1].text = new_content
                # Remove red - copy color from label run or use theme
                if para.runs[0].font.color and para.runs[0].font.color.rgb:
                    try:
                        para.runs[1].font.color.rgb = para.runs[0].font.color.rgb
                    except:
                        para.runs[1].font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                else:
                    para.runs[1].font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                updated = True
    
    return updated


def replace_in_table(table, replacements: Dict[str, str]) -> bool:
    """Replace text in table cells."""
    replaced = False
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in replacements.items():
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                replaced = True
    return replaced


def populate_site_profile_table(table, profile_data) -> bool:
    """
    Populate the Site Profile table with data from SiteProfileData.
    """
    from pptx.util import Pt
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from copy import deepcopy
    
    print("[DEBUG] populate_site_profile_table STARTED")
    
    # Handle duck typing or dict
    if isinstance(profile_data, dict):
        print("[DEBUG] profile_data is dict, converting")
        profile_data = SiteProfileData.from_dict(profile_data)
    
    try:
        descriptions = profile_data.to_description_dict()
        print(f"[DEBUG] Generated {len(descriptions)} descriptions")
        print(f"[DEBUG] Description keys: {list(descriptions.keys())}")
    except Exception as e:
        print(f"[DEBUG] Failed to generate descriptions: {e}")
        return False

    replaced = False
    
    # Create ordered mapping from row item names to description keys
    row_mapping = [
        ('wastewater', 'Wastewater'),
        ('wetlands', 'Wetlands and Jurisdictional Water'),
        ('location', 'Location'),
        ('ownership', 'Ownership & Asking Price'),
        ('size', 'Size / Shape / Dimensions'),
        ('zoning', 'Zoning'),
        ('timing', 'Timing'),
        ('geotechnical', 'Geotechnical and Topography'),
        ('environmental', 'Environmental, Ecological, Archeological'),
        ('disaster', 'Disaster'),
        ('easements', 'Easements'),
        ('electricity', 'Electricity'),
        ('water', 'Water'),
        ('telecom', 'Telecom'),
        ('gas', 'Gas'),
        ('transportation', 'Transportation'),
        ('labor', 'Labor'),
    ]
    
    for row_idx, row in enumerate(table.rows):
        cells = list(row.cells)
        if len(cells) >= 4:
            item_text = cells[0].text.strip().lower()
            print(f"[DEBUG] Row {row_idx}: '{item_text}'")
            
            # Find matching description
            description = None
            for key_fragment, desc_key in row_mapping:
                if key_fragment in item_text:
                    description = descriptions.get(desc_key)
                    print(f"[DEBUG]   Matched '{key_fragment}' -> '{desc_key}'")
                    print(f"[DEBUG]   Description: {description[:50]}..." if description else "   Description: None")
                    break
            
            if description:
                desc_cell = cells[3]
                if desc_cell.text_frame:
                    paragraphs = list(desc_cell.text_frame.paragraphs)
                    
                    if not paragraphs:
                        continue
                    
                    para = paragraphs[0]
                    
                    # Get font properties from first run as template
                    template_font_size = None
                    template_font_name = None
                    if para.runs:
                        first_run = para.runs[0]
                        template_font_size = first_run.font.size
                        template_font_name = first_run.font.name
                    
                    # Clear all existing runs from first paragraph
                    for run in list(para.runs):
                        run._r.getparent().remove(run._r)
                    
                    # Clear any additional paragraphs (they may have red TBD from template)
                    # We need to access the underlying XML to remove extra paragraphs
                    txBody = desc_cell.text_frame._txBody
                    p_elements = txBody.findall(qn('a:p'))
                    # Keep only the first paragraph, remove the rest
                    for p_elem in p_elements[1:]:
                        txBody.remove(p_elem)
                    
                    # Parse the description and create properly formatted runs
                    # Split by periods to handle multiple sentences/fields
                    segments = description.split('. ')
                    
                    for seg_idx, segment in enumerate(segments):
                        if not segment.strip():
                            continue
                        
                        # Add period separator between segments (except first)
                        if seg_idx > 0:
                            sep_run = para.add_run()
                            sep_run.text = ". "
                            if template_font_size:
                                sep_run.font.size = template_font_size
                            if template_font_name:
                                sep_run.font.name = template_font_name
                            sep_run.font.bold = False
                            # Use theme color (dark text)
                            sep_run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                        
                        # Check if segment has a label (contains colon)
                        if ':' in segment:
                            parts = segment.split(':', 1)
                            label = parts[0].strip() + ':'
                            content = parts[1].strip() if len(parts) > 1 else ''
                            
                            # Add bold label
                            label_run = para.add_run()
                            label_run.text = label + ' '
                            label_run.font.bold = True
                            if template_font_size:
                                label_run.font.size = template_font_size
                            if template_font_name:
                                label_run.font.name = template_font_name
                            label_run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                            
                            # Add non-bold content
                            if content:
                                content_run = para.add_run()
                                content_run.text = content
                                content_run.font.bold = False
                                if template_font_size:
                                    content_run.font.size = template_font_size
                                if template_font_name:
                                    content_run.font.name = template_font_name
                                content_run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                        else:
                            # No colon - just add as regular text
                            text_run = para.add_run()
                            text_run.text = segment
                            text_run.font.bold = False
                            if template_font_size:
                                text_run.font.size = template_font_size
                            if template_font_name:
                                text_run.font.name = template_font_name
                            text_run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
                    
                    replaced = True
    
    return replaced


def build_replacements(site_data: Dict, config: ExportConfig) -> Dict[str, str]:
    """Build replacement dictionary."""
    target_mw = site_data.get('target_mw', 0)
    mw_display = f"{target_mw/1000:.1f}GW" if target_mw >= 1000 else f"{target_mw:.0f}MW"
    
    # Get profile if available
    profile = site_data.get('profile')
    if isinstance(profile, dict):
        profile = SiteProfileData.from_dict(profile)

    replacements = {
        'SITE NAME': site_data.get('name', 'Site Name'),
        'Site Name': site_data.get('name', 'Site Name'),  # Also in left panel
        '[STATE]': site_data.get('state', 'State'),
        'Oklahoma': site_data.get('state', 'State'),  # State name in left panel
        'December 2, 2025': datetime.now().strftime('%B %d, %Y'),
        '[Coordinates linked]': f"{site_data.get('latitude', 'N/A')}, {site_data.get('longitude', 'N/A')}",
        '1GW+': mw_display,
        'Estimated 1GW+': f"Estimated {mw_display}",
        'NAME': config.contact_name or site_data.get('contact_name', 'Contact Name'),
        'TITLE': config.contact_title or site_data.get('contact_title', 'Title'),
        'Phone': config.contact_phone or site_data.get('contact_phone', 'Phone'),
        'Email': config.contact_email or site_data.get('contact_email', 'email@jll.com'),
        # Last Edited date
        '12.1.2025': datetime.now().strftime('%m.%d.%Y'),
        'Last Edited: 12.1.2025': f"Last Edited: {datetime.now().strftime('%m.%d.%Y')}",
    }

    if site_data.get('utility'):
        replacements['OG&E 345kV line'] = site_data['utility']
    if site_data.get('total_acres') or site_data.get('acreage'):
        acres = site_data.get('total_acres') or site_data.get('acreage')
        replacements['1,250 acres'] = f"{acres:,.0f} acres"
    if site_data.get('water_capacity'):
        replacements['3M GDP capacity'] = site_data['water_capacity']
    
    # Overview, Observation, Outstanding are handled by update_overview_textbox()
    # But we still need the overview replacement for the find_and_replace fallback
    if profile and profile.overview:
        replacements['1,250 acre site with significant growth opportunity located between Tulsa and Oklahoma City.'] = profile.overview
    elif site_data.get('overview'):
        replacements['1,250 acre site with significant growth opportunity located between Tulsa and Oklahoma City.'] = site_data['overview']
    
    # Also handle if overview/observation/outstanding are directly in site_data
    if site_data.get('observation') and site_data['observation'] != 'to be completed':
        replacements['Observation:  to be completed'] = f"Observation:  {site_data['observation']}"
    if site_data.get('outstanding') and site_data['outstanding'] != 'to be completed':
        replacements['Outstanding:  to be completed'] = f"Outstanding:  {site_data['outstanding']}"

    return replacements


# =============================================================================
# MAIN EXPORT
# =============================================================================


def add_critical_path_text(slide, phases):
    """Add Critical Path text to slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(6.0)
    height = Inches(5.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Critical Path to Power"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 43, 74) # Dark Blue
    
    for phase in phases:
        p = tf.add_paragraph()
        p.text = f"Phase {phase.phase_num}: {phase.target_mw:.0f} MW @ {phase.voltage_kv} kV"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 43, 74)
        p.space_before = Pt(12)
        
        p = tf.add_paragraph()
        p.text = f"Target: {phase.target_online}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(102, 102, 102) # Medium Gray
        
        items = [
            ('Screening Study', phase.screening_study),
            ('Contract Study', phase.contract_study),
            ('Letter of Agreement', phase.letter_of_agreement),
            ('Energy Contract', phase.energy_contract),
        ]
        for label, status in items:
            p = tf.add_paragraph()
            if status.lower() in ['complete', 'executed']:
                symbol = '✓'
            elif status.lower() in ['drafted', 'initiated', 'in_progress', 'in progress']:
                symbol = '○'
            else:
                symbol = '□'
            
            p.text = f"{symbol} {label}: {status}"
            p.font.size = Pt(11)
            p.level = 1
            p.space_before = Pt(3)


def add_score_breakdown_text(slide, scores):
    """Add Score Breakdown text to slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    left = Inches(6.8)
    top = Inches(1.5)
    width = Inches(6.0)
    height = Inches(5.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Score Breakdown"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 43, 74)
    
    p = tf.add_paragraph()
    p.text = f"Overall Score: {scores.overall_score:.0f}/100"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_before = Pt(12)
    if scores.overall_score >= 70:
        p.font.color.rgb = RGBColor(46, 125, 50) # Green
    elif scores.overall_score >= 50:
        p.font.color.rgb = RGBColor(249, 168, 37) # Amber
    else:
        p.font.color.rgb = RGBColor(227, 24, 55) # Red
        
    sub_scores = [
        ('Power Pathway', scores.power_pathway, 0.30),
        ('Site Specific', scores.site_specific, 0.10),
        ('Execution', scores.execution, 0.20),
        ('Relationships', scores.relationship_capital, 0.35),
        ('Financial', scores.financial, 0.05),
    ]
    
    for name, score, weight in sub_scores:
        p = tf.add_paragraph()
        p.text = f"{name}: {score:.0f} (Weight: {weight*100:.0f}%)"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_before = Pt(10)


def add_market_text(slide, market_data):
    """Add Market Analysis text quadrants."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    # Top Right: Competitive Landscape
    left, top, width, height = Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Competitive Landscape"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 43, 74)
    
    existing_mw = market_data.get('existing_dc_mw', 0)
    fiber = market_data.get('fiber_density', 'medium')
    hyperscalers = market_data.get('hyperscaler_presence', [])
    
    p = tf.add_paragraph()
    p.text = f"Existing DC Capacity: {existing_mw:,} MW"
    p.font.size = Pt(12)
    p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = f"Fiber Density: {fiber.title()}"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Major Operators:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.space_before = Pt(6)
    
    if hyperscalers:
        for hs in hyperscalers[:4]:
            p = tf.add_paragraph()
            p.text = f"• {hs}"
            p.font.size = Pt(11)
            p.level = 1
    else:
        p = tf.add_paragraph()
        p.text = "• None identified"
        p.font.size = Pt(11)
        p.level = 1

    # Bottom Left: ISO/Utility
    left, top = Inches(0.5), Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ISO / Utility Profile"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 43, 74)
    
    iso = market_data.get('primary_iso', 'N/A')
    utility = market_data.get('utility_name', 'N/A')
    queue = market_data.get('avg_queue_time_months', 'N/A')
    
    items = [
        f"ISO: {iso}",
        f"Utility: {utility}",
        f"Avg Queue Time: {queue} months",
        f"Regulatory: {market_data.get('regulatory_structure', 'N/A').title()}"
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.space_before = Pt(6)

    # Bottom Right: SWOT
    left, top = Inches(6.8), Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "State SWOT Summary"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 43, 74)
    
    strengths = market_data.get('swot', {}).get('strengths', [])
    weaknesses = market_data.get('swot', {}).get('weaknesses', [])
    
    p = tf.add_paragraph()
    p.text = "Strengths:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(46, 125, 50)
    p.space_before = Pt(6)
    
    for s in strengths[:2]:
        p = tf.add_paragraph()
        p.text = f"• {s}"
        p.font.size = Pt(11)
        p.level = 1
        
    p = tf.add_paragraph()
    p.text = "Challenges:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(249, 168, 37)
    p.space_before = Pt(6)
    
    for w in weaknesses[:2]:
        p = tf.add_paragraph()
        p.text = f"• {w}"
        p.font.size = Pt(11)
        p.level = 1


def convert_phase_data(phase_dict: Dict) -> Dict:
    """Convert Google Sheets phase data format to PhaseData format."""
    # Google Sheets uses: mw, screening_status, contract_study_status, loa_status, energy_contract_status, target_date, voltage
    # PhaseData expects: target_mw, screening_study, contract_study, letter_of_agreement, energy_contract, target_online, voltage_kv
    
    converted = {}
    
    # Map MW
    if 'mw' in phase_dict:
        converted['target_mw'] = float(phase_dict['mw'])
    elif 'target_mw' in phase_dict:
        converted['target_mw'] = float(phase_dict['target_mw'])
    else:
        converted['target_mw'] = 0.0
    
    # Map voltage
    if 'voltage' in phase_dict:
        try:
            converted['voltage_kv'] = int(phase_dict['voltage'])
        except (ValueError, TypeError):
            converted['voltage_kv'] = 138
    elif 'voltage_kv' in phase_dict:
        converted['voltage_kv'] = int(phase_dict['voltage_kv'])
    else:
        converted['voltage_kv'] = 138
    
    # Map target date/online
    if 'target_date' in phase_dict:
        converted['target_online'] = phase_dict['target_date']
    elif 'target_online' in phase_dict:
        converted['target_online'] = phase_dict['target_online']
    else:
        converted['target_online'] = '2028-01-01'
    
    # Map study statuses
    if 'screening_status' in phase_dict:
        converted['screening_study'] = phase_dict['screening_status']
    elif 'screening_study' in phase_dict:
        converted['screening_study'] = phase_dict['screening_study']
    else:
        converted['screening_study'] = 'Not Started'
    
    if 'contract_study_status' in phase_dict:
        converted['contract_study'] = phase_dict['contract_study_status']
    elif 'contract_study' in phase_dict:
        converted['contract_study'] = phase_dict['contract_study']
    else:
        converted['contract_study'] = 'Not Started'
    
    if 'loa_status' in phase_dict:
        converted['letter_of_agreement'] = phase_dict['loa_status']
    elif 'letter_of_agreement' in phase_dict:
        converted['letter_of_agreement'] = phase_dict['letter_of_agreement']
    else:
        converted['letter_of_agreement'] = 'Not Started'
    
    if 'energy_contract_status' in phase_dict:
        converted['energy_contract'] = phase_dict['energy_contract_status']
    elif 'energy_contract' in phase_dict:
        converted['energy_contract'] = phase_dict['energy_contract']
    else:
        converted['energy_contract'] = 'Not Started'
    
    # Optional fields
    if 'service_type' in phase_dict:
        converted['transmission_type'] = phase_dict['service_type']
    elif 'transmission_type' in phase_dict:
        converted['transmission_type'] = phase_dict['transmission_type']
    else:
        converted['transmission_type'] = ''
    
    if 'substation_status' in phase_dict:
        converted['substation_type'] = phase_dict['substation_status']
    elif 'substation_type' in phase_dict:
        converted['substation_type'] = phase_dict['substation_type']
    else:
        converted['substation_type'] = ''
    
    if 'trans_dist' in phase_dict:
        try:
            converted['distance_to_transmission'] = float(phase_dict['trans_dist'])
        except (ValueError, TypeError):
            converted['distance_to_transmission'] = 0.0
    elif 'distance_to_transmission' in phase_dict:
        converted['distance_to_transmission'] = float(phase_dict['distance_to_transmission'])
    else:
        converted['distance_to_transmission'] = 0.0
    
    # Add phase number if missing
    if 'phase_num' not in converted and 'phase_num' in phase_dict:
        converted['phase_num'] = int(phase_dict['phase_num'])
    elif 'phase_num' not in converted:
        converted['phase_num'] = 1
    
    return converted


def populate_slide(slide, site_data, profile_data, replacements, config, slide_type=None):
    """
    Populate a single slide with site data.
    Shared logic between single-site export and portfolio export.
    """
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    print(f"[DEBUG] Populating slide (Type: {slide_type})")

    # 1. Text and Table Replacements
    for shape in slide.shapes:
        if shape.has_table:
            row_count = len(shape.table.rows)
            # Check if this is the main Site Profile table (18 rows)
            if row_count >= 17 and (slide_type == 'site_profile' or slide_type is None):
                if profile_data:
                    populate_site_profile_table(shape.table, profile_data)
                else:
                    replace_in_table(shape.table, replacements)
            else:
                replace_in_table(shape.table, replacements)
        elif shape.has_text_frame:
            # Check if this is the Overview/Observation/Outstanding text box
            shape_text = shape.text_frame.text.lower()
            if 'overview' in shape_text and 'observation' in shape_text:
                update_overview_textbox(shape, site_data, profile_data)
            else:
                find_and_replace_text(shape, replacements)
                
        # Recursive check for groups
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for member in shape.shapes:
                if member.has_text_frame:
                    find_and_replace_text(member, replacements)

    # 2. Site Boundary & Topography Placeholders
    if slide_type in ['site_boundary', 'topography']:
        # Find pictures to replace
        pics_to_replace = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics_to_replace.append(shape)
        
        # Replace them
        for pic in pics_to_replace:
            left, top, width, height = pic.left, pic.top, pic.width, pic.height
            # Remove picture
            sp = pic._element
            sp.getparent().remove(sp)
            
            # Add placeholder
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
            shape.line.color.rgb = RGBColor(200, 200, 200)
            
            # Add label
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = "Site Boundary Map" if slide_type == 'site_boundary' else "Topography Map"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.font.size = Pt(18)

    # 3. Site Profile Specifics (State Silhouette, Map Placeholders)
    if slide_type == 'site_profile':
        shapes_to_replace = []
        state_shape = None
        
        for shape in slide.shapes:
            # Check for state silhouette on the left (approx < 4 inches)
            if shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE] and shape.left < Inches(4):
                # Heuristic: It's likely the state shape if it's in the top-left quadrant
                if shape.top < Inches(4):
                    state_shape = shape

            # Check for map images on the right (approx > 9 inches)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.left > Inches(9):
                shapes_to_replace.append(shape)
        
        # Handle state silhouette
        current_state = site_data.get('state', '').upper()
        if state_shape and current_state != 'OK':
            # Hide the Oklahoma shape if state is not OK
            sp_element = state_shape.element
            sp_element.getparent().remove(sp_element)

        # Replace map images with placeholders
        for shape in shapes_to_replace:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # Remove original
            sp_element = shape.element
            sp_element.getparent().remove(sp_element)
            
            # Add placeholder rectangle
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(240, 240, 240) # Light gray
            rect.line.color.rgb = RGBColor(200, 200, 200)
            
            # Add text
            tf = rect.text_frame
            tf.text = "Map Placeholder"
            # Determine type based on vertical position
            if top < Inches(2.5):
                tf.text = "Location Map"
            elif top < Inches(5):
                tf.text = "Site Map"
            else:
                tf.text = "Plot Map"
            
            p = tf.paragraphs[0]
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.alignment = PP_ALIGN.CENTER


def export_site_to_pptx(
    site_data: Dict,
    template_path: str,
    output_path: str,
    config: ExportConfig = None,
) -> str:
    """Export site data to PowerPoint."""
    print(f"[DEBUG] export_site_to_pptx called with template_path='{template_path}', output_path='{output_path}'")
    
    if not template_path:
        raise ValueError("template_path is required and cannot be None or empty")
    if not output_path:
        raise ValueError("output_path is required and cannot be None or empty")
        
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("python-pptx required")

    config = config or ExportConfig()
    prs = Presentation(template_path)
    replacements = build_replacements(site_data, config)
    
    # Check if we have structured profile data
    profile_data = None
    if 'profile' in site_data:
        p = site_data['profile']
        # Use duck typing instead of isinstance to handle module reloads
        if hasattr(p, 'overview') and hasattr(p, 'to_description_dict'):
            profile_data = p
        elif isinstance(p, dict):
            profile_data = SiteProfileData.from_dict(p)

    # Process existing slides
    for slide_idx, slide in enumerate(prs.slides):
        print(f"[DEBUG] Processing slide {slide_idx}")
        
        # Determine slide type based on index
        slide_type = None
        if slide_idx == 1:
            slide_type = 'site_profile'
        elif slide_idx == 2:
            slide_type = 'site_boundary'
        elif slide_idx == 3:
            slide_type = 'topography'
            
        populate_slide(slide, site_data, profile_data, replacements, config, slide_type)

    # Find blank layout
    blank_layout = None
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            blank_layout = layout
            break
    
    # Fallback to last layout if no blank found
    blank_layout = blank_layout or prs.slide_layouts[-1]
    
    # Clean "Thank You" text from the selected layout
    # This ensures new slides don't inherit the "Thank You" background
    try:
        shapes_to_remove = []
        for shape in blank_layout.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                if "thank you" in text or "questions?" in text:
                    shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
    except Exception as e:
        print(f"[WARNING] Failed to clean layout: {e}")

    # Get trajectory and phases
    trajectory = None
    if config.include_capacity_trajectory and MATPLOTLIB_AVAILABLE:
        # Try 'capacity_trajectory' first, then 'schedule'
        traj_data = site_data.get('capacity_trajectory', site_data.get('schedule', {}))
        print(f"[DEBUG] Trajectory data found: {bool(traj_data)}")
        if traj_data:
            print(f"[DEBUG] Trajectory data keys: {list(traj_data.keys())}")
        
        trajectory = (CapacityTrajectory.from_dict(traj_data) if traj_data else
                      CapacityTrajectory.generate_default(
                          site_data.get('target_mw', 600),
                          site_data.get('phase1_mw'),
                          site_data.get('start_year', 2028)))



    phases = []
    phase_data = site_data.get('phases', [])
    if phase_data:
        for idx, pd in enumerate(phase_data, 1):
            if isinstance(pd, dict):
                # Convert Google Sheets format to PhaseData format
                converted_pd = convert_phase_data(pd)
                # Ensure phase_num is set
                if 'phase_num' not in converted_pd:
                    converted_pd['phase_num'] = idx
                phases.append(PhaseData(**converted_pd))
            else:
                phases.append(pd)
    else:
        target_mw = site_data.get('target_mw', 600)
        phase1_mw = site_data.get('phase1_mw', min(100, target_mw * 0.2))
        phases = [
            PhaseData(1, phase1_mw, 138, site_data.get('phase1_online', '2028-01-01'),
                      site_data.get('screening_study', 'Complete'),
                      site_data.get('contract_study', 'In Progress'),
                      site_data.get('loa_status', 'Not Started'),
                      site_data.get('energy_contract', 'Not Started')),
            PhaseData(2, target_mw, 345, site_data.get('phase2_online', '2029-01-01'),
                      site_data.get('phase2_screening', 'Complete'),
                      site_data.get('phase2_contract_study', 'Not Started')),
        ]

    # ADD SLIDE: Capacity Trajectory (WHITE background)
    if config.include_capacity_trajectory and trajectory:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set white background explicitly
        set_slide_background_white(slide, Inches, RGBColor)
        
        add_header_bar(slide, f"{site_data.get('name', 'Site')}: Capacity Trajectory",
                      Inches, Pt, RGBColor)

        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            chart_path = tmp.name
        print(f"[DEBUG] Generating trajectory chart at {chart_path}")
        generate_capacity_trajectory_chart(trajectory, site_data.get('name', 'Site'),
                                           chart_path, phases, subtitle=config.chart_subtitle)
        slide.shapes.add_picture(chart_path, Inches(0.4), Inches(0.8), width=Inches(12.5))
        add_footer(slide, 5, Inches, Pt, RGBColor)
        os.unlink(chart_path)

    # ADD SLIDE: Infrastructure & Critical Path (WHITE background)
    if config.include_infrastructure and MATPLOTLIB_AVAILABLE:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set white background explicitly
        set_slide_background_white(slide, Inches, RGBColor)
        
        add_header_bar(slide, "Infrastructure & Critical Path", Inches, Pt, RGBColor)

        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            infra_path = tmp.name
        print(f"[DEBUG] Generating infra chart at {infra_path}")
        # Generate Infrastructure Chart (Right side)
        generate_critical_path_chart(phases, site_data, infra_path, width=6, height=5)
        slide.shapes.add_picture(infra_path, Inches(6.8), Inches(1.5), width=Inches(6.0))
        
        # Add Critical Path Text (Left side)
        add_critical_path_text(slide, phases)

        # Add risks/opportunities summary at bottom
        risks = site_data.get('risks', [])
        opportunities = site_data.get('opportunities', [])
        if risks or opportunities:
            details_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8))
            tf = details_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            summary_parts = []
            if risks:
                summary_parts.append(f"Key Risks: {', '.join(risks[:3])}")
            if opportunities:
                summary_parts.append(f"Opportunities: {', '.join(opportunities[:3])}")
            p.text = "  |  ".join(summary_parts)
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor.from_string(JLL_COLORS['medium_gray'][1:])

        add_footer(slide, 6, Inches, Pt, RGBColor)
        os.unlink(infra_path)

    # ADD SLIDE: Score Analysis (WHITE background)
    if config.include_score_analysis and MATPLOTLIB_AVAILABLE:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background_white(slide, Inches, RGBColor)
        add_header_bar(slide, "Site Score Analysis", Inches, Pt, RGBColor)
        
        # Get or create scores
        scores_data = site_data.get('scores', {})
        scores = ScoreAnalysis(
            overall_score=scores_data.get('overall', site_data.get('overall_score', 65)),
            power_pathway=scores_data.get('power_pathway', site_data.get('power_score', 70)),
            site_specific=scores_data.get('site_specific', site_data.get('site_score', 60)),
            execution=scores_data.get('execution', site_data.get('execution_score', 55)),
            relationship_capital=scores_data.get('relationship_capital', site_data.get('relationship_score', 50)),
            financial=scores_data.get('financial', site_data.get('financial_score', 75)),
            land_score=scores_data.get('land', 0),
            queue_score=scores_data.get('queue', 0),
            water_score=scores_data.get('water', 0),
            fiber_score=scores_data.get('fiber', 0),
        )
        
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            score_path = tmp.name
        print(f"[DEBUG] Generating score chart at {score_path}")
        generate_score_radar_chart(scores, site_data.get('name', 'Site'), score_path, width=6, height=6)
        slide.shapes.add_picture(score_path, Inches(0.5), Inches(1.5), width=Inches(6.0))
        
        # Add Score Breakdown Text (Right side)
        add_score_breakdown_text(slide, scores)
        
        add_footer(slide, 7, Inches, Pt, RGBColor)
        os.unlink(score_path)

    # ADD SLIDE: Market Analysis (WHITE background)
    if config.include_market_analysis and MATPLOTLIB_AVAILABLE:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background_white(slide, Inches, RGBColor)
        add_header_bar(slide, "Market Analysis", Inches, Pt, RGBColor)
        
        # Build market data from site_data and state analysis
        state_code = site_data.get('state_code', site_data.get('state', 'OK')[:2].upper())
        
        # Get market_analysis if provided, otherwise build from available data
        market_raw = site_data.get('market_analysis', {})
        
        # Build comparison states data
        comparison_states = market_raw.get('comparison_states', {})
        if not comparison_states:
            # Default comparison states based on common alternatives
            default_comparisons = {
                'OK': {'TX': {}, 'GA': {}, 'OH': {}},
                'TX': {'OK': {}, 'GA': {}, 'AZ': {}},
                'GA': {'TX': {}, 'VA': {}, 'OH': {}},
                'VA': {'GA': {}, 'OH': {}, 'TX': {}},
                'OH': {'IN': {}, 'GA': {}, 'TX': {}},
                'IN': {'IN': {}, 'GA': {}, 'TX': {}},
                'WY': {'IN': {}, 'GA': {}, 'TX': {}},
                'AZ': {'IN': {}, 'GA': {}, 'TX': {}},
            }
            comp_codes = list(default_comparisons.get(state_code, {'TX': {}, 'GA': {}}).keys())[:3]
            
            # Default state data (from state_analysis framework)
            state_defaults = {
                'OK': {'name': 'Oklahoma', 'queue_months': 30, 'rate': 0.055, 'score': 88, 'tier': 1},
                'TX': {'name': 'Texas', 'queue_months': 36, 'rate': 0.065, 'score': 80, 'tier': 1},
                'GA': {'name': 'Georgia', 'queue_months': 42, 'rate': 0.072, 'score': 70, 'tier': 2},
                'VA': {'name': 'Virginia', 'queue_months': 48, 'rate': 0.078, 'score': 58, 'tier': 3},
                'OH': {'name': 'Ohio', 'queue_months': 36, 'rate': 0.068, 'score': 70, 'tier': 2},
                'IN': {'name': 'Indiana', 'queue_months': 32, 'rate': 0.062, 'score': 72, 'tier': 2},
                'WY': {'name': 'Wyoming', 'queue_months': 28, 'rate': 0.048, 'score': 82, 'tier': 1},
                'AZ': {'name': 'Arizona', 'queue_months': 38, 'rate': 0.070, 'score': 52, 'tier': 3},
            }
            for code in comp_codes:
                if code in state_defaults:
                    comparison_states[code] = state_defaults[code]
        
        # Get state defaults for site's state
        site_state_defaults = {
            'OK': {'iso': 'SPP', 'reg': 'regulated', 'queue': 30, 'rate': 0.055, 'renew': 42,
                   'score': 88, 'tier': 1, 'dc_mw': 500, 'fiber': 'medium',
                   'hyperscalers': ['Google (Pryor)', 'Meta (announced)'],
                   'strengths': ['Pro-business PSC', 'Low power costs', 'SPP wholesale market'],
                   'weaknesses': ['Limited DC ecosystem', 'Water constraints'],
                   'opportunities': ['Hyperscaler expansion', 'Tulsa hub growth'],
                   'threats': ['Grid congestion', 'Water rights competition'],
                   'incentives': ['Sales tax exemption', 'Property tax abatement', 'Quality Jobs Program']},
            'TX': {'iso': 'ERCOT', 'reg': 'deregulated', 'queue': 36, 'rate': 0.065, 'renew': 35,
                   'score': 80, 'tier': 1, 'dc_mw': 3000, 'fiber': 'high',
                   'hyperscalers': ['Google', 'Microsoft', 'Meta', 'AWS', 'Oracle'],
                   'strengths': ['No state income tax', 'Massive ecosystem', 'ERCOT flexibility'],
                   'weaknesses': ['Grid reliability', 'Water stress', 'Queue backlog'],
                   'opportunities': ['Continued growth', 'West Texas expansion'],
                   'threats': ['Grid instability', 'Water availability'],
                   'incentives': ['Chapter 313 replacement', 'Property tax limits']},
        }
        
        defaults = site_state_defaults.get(state_code, site_state_defaults.get('OK'))
        
        market_data = {
            'state_code': state_code,
            'state_name': market_raw.get('state_name', site_data.get('state', defaults.get('name', 'State'))),
            'primary_iso': market_raw.get('primary_iso', site_data.get('iso', defaults.get('iso', 'SPP'))),
            'regulatory_structure': market_raw.get('regulatory_structure', defaults.get('reg', 'regulated')),
            'utility_name': market_raw.get('utility_name', site_data.get('utility', 'Utility')),
            'avg_queue_time_months': market_raw.get('avg_queue_time_months', defaults.get('queue', 36)),
            'avg_industrial_rate': market_raw.get('avg_industrial_rate', defaults.get('rate', 0.06)),
            'renewable_percentage': market_raw.get('renewable_percentage', defaults.get('renew', 30)),
            'overall_score': market_raw.get('overall_score', defaults.get('score', 70)),
            'tier': market_raw.get('tier', defaults.get('tier', 2)),
            'existing_dc_mw': market_raw.get('existing_dc_mw', defaults.get('dc_mw', 500)),
            'hyperscaler_presence': market_raw.get('hyperscaler_presence', defaults.get('hyperscalers', [])),
            'fiber_density': market_raw.get('fiber_density', defaults.get('fiber', 'medium')),
            'swot': {
                'strengths': market_raw.get('strengths', defaults.get('strengths', [])),
                'weaknesses': market_raw.get('weaknesses', defaults.get('weaknesses', [])),
                'opportunities': market_raw.get('opportunities', defaults.get('opportunities', [])),
                'threats': market_raw.get('threats', defaults.get('threats', [])),
            },
            'comparison_states': comparison_states,
            'incentives': market_raw.get('incentives', defaults.get('incentives', [])),
        }
        
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            market_path = tmp.name
            
        print(f"[DEBUG] Generating market chart at {market_path}")
        # Generate Market Chart (Top Left)
        generate_market_analysis_chart(market_data, site_data.get('name', 'Site'), market_path, width=6, height=3.5)
        slide.shapes.add_picture(market_path, Inches(0.5), Inches(1.5), width=Inches(6.0))
        
        # Add Market Text Quadrants
        add_market_text(slide, market_data)
        
        add_footer(slide, 8, Inches, Pt, RGBColor)
        os.unlink(market_path)

    # REORDER: Move Thank You to end
    def move_slide_to_end(prs, slide_index):
        slides = list(prs.slides._sldIdLst)
        if slide_index < len(slides):
            slide_id = slides[slide_index]
            prs.slides._sldIdLst.remove(slide_id)
            prs.slides._sldIdLst.append(slide_id)

    move_slide_to_end(prs, 4)  # Original Thank You position

    prs.save(output_path)
    return output_path


def export_multiple_sites(sites: List[Dict], template_path: str,
                          output_dir: str, config: ExportConfig = None) -> List[str]:
    """Export multiple sites."""
    config = config or ExportConfig()
    os.makedirs(output_dir, exist_ok=True)
    output_paths = []
    for site in sites:
        site_id = site.get('site_id', 'site')
        output_path = os.path.join(output_dir, f"{site_id}_profile.pptx")
        try:
            export_site_to_pptx(site, template_path, output_path, config)
            output_paths.append(output_path)
        except Exception as e:
            print(f"Error exporting {site_id}: {e}")
    return output_paths


def analyze_template(template_path: str) -> Dict:
    """Analyze template structure."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx required")

    prs = Presentation(template_path)
    analysis = {'slide_count': len(prs.slides), 'slides': []}
    for i, slide in enumerate(prs.slides):
        slide_info = {'index': i, 'shapes': [], 'tables': []}
        for shape in slide.shapes:
            if shape.has_table:
                slide_info['tables'].append({
                    'rows': len(shape.table.rows),
                    'cols': len(shape.table.columns)
                })
            elif shape.has_text_frame and shape.text_frame.text.strip():
                slide_info['shapes'].append({'text': shape.text_frame.text[:100]})
        analysis['slides'].append(slide_info)
    return analysis


def create_default_template(output_path: str):
    """Create a default PowerPoint template."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError("python-pptx required")
    
    prs = Presentation()
    
    # Slide 0: Title Slide (Dark Blue)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 43, 74) # Dark Blue
    
    # Slide 1: Site Profile (White)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide 2: Site Boundary (White)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide 3: Topography (White)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide 4: Capacity Trajectory (White)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide 5: Infrastructure & Critical Path (White)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Slide 6: Thank You (Dark Blue)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 43, 74)
    
    prs.save(output_path)


__all__ = [
    'CapacityTrajectory', 'PhaseData', 'ScoreAnalysis', 'RiskOpportunity', 
    'SiteProfileData', 'MarketAnalysis',
    'ExportConfig', 'export_site_to_pptx', 'export_multiple_sites',
    'generate_capacity_trajectory_chart', 'generate_critical_path_chart',
    'generate_score_radar_chart', 'generate_score_summary_chart',
    'generate_market_analysis_chart', 'analyze_template', 'JLL_COLORS',
    'MATPLOTLIB_AVAILABLE', 'TEMPLATE_VERSION', 'create_default_template',
]
