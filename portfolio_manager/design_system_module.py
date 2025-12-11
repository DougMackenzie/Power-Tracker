"""
Design System Learning Module
=============================
Drop-in module for learning presentation styles from existing PPTX files
and enforcing consistent styling when generating new presentations.

Integration:
1. Copy this file to your portfolio_manager folder
2. Import: from design_system_module import DesignSystemAnalyzer, StyleEnforcer, DesignSystem
3. Train: design_system = DesignSystemAnalyzer().analyze(["file1.pptx", "file2.pptx"])
4. Use: enforcer = StyleEnforcer(design_system) in your PPTX generation code

Author: Presentation Builder
Version: 1.0.0
"""

import os
import json
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from dataclasses import dataclass, field, asdict
from datetime import datetime
from typing import Dict, List, Optional, Tuple, Any
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# =============================================================================
# DATA STRUCTURES
# =============================================================================

@dataclass
class ColorValue:
    """A color with hex value."""
    hex: str
    usage: str = ""  # "primary", "secondary", "text", etc.
    
    def __post_init__(self):
        # Normalize hex
        if not self.hex.startswith("#"):
            self.hex = f"#{self.hex}"
        self.hex = self.hex.upper()
    
    def to_rgb(self) -> RGBColor:
        """Convert to python-pptx RGBColor."""
        h = self.hex.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    
    def to_dict(self) -> dict:
        return {"hex": self.hex, "usage": self.usage}
    
    @classmethod
    def from_dict(cls, d: dict) -> "ColorValue":
        return cls(hex=d["hex"], usage=d.get("usage", ""))


@dataclass
class TypographyStyle:
    """A typography style specification."""
    font_name: str
    font_size: float  # in points
    bold: bool = False
    italic: bool = False
    color: str = "text_primary"  # Color name reference
    
    def to_dict(self) -> dict:
        return asdict(self)
    
    @classmethod
    def from_dict(cls, d: dict) -> "TypographyStyle":
        return cls(**d)


@dataclass
class DesignSystem:
    """
    Complete design system learned from presentations.
    
    This ensures consistent styling across all generated presentations.
    """
    # Metadata
    name: str = "Organization"
    version: str = "1.0.0"
    created: str = field(default_factory=lambda: datetime.now().isoformat())
    
    # Colors
    primary_color: ColorValue = field(default_factory=lambda: ColorValue("#1A2B4A"))
    secondary_color: ColorValue = field(default_factory=lambda: ColorValue("#666666"))
    accent_color: Optional[ColorValue] = None
    text_primary: ColorValue = field(default_factory=lambda: ColorValue("#333333"))
    text_secondary: ColorValue = field(default_factory=lambda: ColorValue("#666666"))
    background: ColorValue = field(default_factory=lambda: ColorValue("#FFFFFF"))
    
    # Additional brand colors
    additional_colors: Dict[str, ColorValue] = field(default_factory=dict)
    
    # Typography
    title_style: TypographyStyle = field(default_factory=lambda: TypographyStyle("Arial", 28, bold=True))
    h1_style: TypographyStyle = field(default_factory=lambda: TypographyStyle("Arial", 22, bold=True))
    h2_style: TypographyStyle = field(default_factory=lambda: TypographyStyle("Arial", 16, bold=True))
    body_style: TypographyStyle = field(default_factory=lambda: TypographyStyle("Arial", 11))
    caption_style: TypographyStyle = field(default_factory=lambda: TypographyStyle("Arial", 9))
    metric_style: TypographyStyle = field(default_factory=lambda: TypographyStyle("Arial", 36, bold=True))
    
    # Table styling
    table_header_bg: str = "primary"  # Reference to color name
    table_header_text: str = "#FFFFFF"
    table_alternate_rows: bool = True
    
    # Spacing (inches)
    margin_left: float = 0.5
    margin_right: float = 0.5
    margin_top: float = 0.5
    margin_bottom: float = 0.5
    element_gap: float = 0.2
    
    # Analysis metadata
    files_analyzed: int = 0
    slides_analyzed: int = 0
    
    def get_color(self, name: str) -> ColorValue:
        """Get a color by name."""
        color_map = {
            "primary": self.primary_color,
            "secondary": self.secondary_color,
            "accent": self.accent_color,
            "text_primary": self.text_primary,
            "text_secondary": self.text_secondary,
            "background": self.background,
        }
        if name in color_map and color_map[name]:
            return color_map[name]
        if name in self.additional_colors:
            return self.additional_colors[name]
        return self.primary_color
    
    def get_typography(self, style_name: str) -> TypographyStyle:
        """Get a typography style by name."""
        style_map = {
            "title": self.title_style,
            "h1": self.h1_style,
            "h2": self.h2_style,
            "body": self.body_style,
            "caption": self.caption_style,
            "metric": self.metric_style,
        }
        return style_map.get(style_name, self.body_style)
    
    def to_dict(self) -> dict:
        """Convert to dictionary for JSON serialization."""
        return {
            "name": self.name,
            "version": self.version,
            "created": self.created,
            "colors": {
                "primary": self.primary_color.to_dict(),
                "secondary": self.secondary_color.to_dict(),
                "accent": self.accent_color.to_dict() if self.accent_color else None,
                "text_primary": self.text_primary.to_dict(),
                "text_secondary": self.text_secondary.to_dict(),
                "background": self.background.to_dict(),
                "additional": {k: v.to_dict() for k, v in self.additional_colors.items()},
            },
            "typography": {
                "title": self.title_style.to_dict(),
                "h1": self.h1_style.to_dict(),
                "h2": self.h2_style.to_dict(),
                "body": self.body_style.to_dict(),
                "caption": self.caption_style.to_dict(),
                "metric": self.metric_style.to_dict(),
            },
            "table": {
                "header_bg": self.table_header_bg,
                "header_text": self.table_header_text,
                "alternate_rows": self.table_alternate_rows,
            },
            "spacing": {
                "margin_left": self.margin_left,
                "margin_right": self.margin_right,
                "margin_top": self.margin_top,
                "margin_bottom": self.margin_bottom,
                "element_gap": self.element_gap,
            },
            "analysis": {
                "files_analyzed": self.files_analyzed,
                "slides_analyzed": self.slides_analyzed,
            }
        }
    
    @classmethod
    def from_dict(cls, d: dict) -> "DesignSystem":
        """Create from dictionary."""
        ds = cls()
        ds.name = d.get("name", "Organization")
        ds.version = d.get("version", "1.0.0")
        ds.created = d.get("created", datetime.now().isoformat())
        
        colors = d.get("colors", {})
        if "primary" in colors:
            ds.primary_color = ColorValue.from_dict(colors["primary"])
        if "secondary" in colors:
            ds.secondary_color = ColorValue.from_dict(colors["secondary"])
        if colors.get("accent"):
            ds.accent_color = ColorValue.from_dict(colors["accent"])
        if "text_primary" in colors:
            ds.text_primary = ColorValue.from_dict(colors["text_primary"])
        if "text_secondary" in colors:
            ds.text_secondary = ColorValue.from_dict(colors["text_secondary"])
        if "background" in colors:
            ds.background = ColorValue.from_dict(colors["background"])
        if "additional" in colors:
            ds.additional_colors = {k: ColorValue.from_dict(v) for k, v in colors["additional"].items()}
        
        typography = d.get("typography", {})
        if "title" in typography:
            ds.title_style = TypographyStyle.from_dict(typography["title"])
        if "h1" in typography:
            ds.h1_style = TypographyStyle.from_dict(typography["h1"])
        if "h2" in typography:
            ds.h2_style = TypographyStyle.from_dict(typography["h2"])
        if "body" in typography:
            ds.body_style = TypographyStyle.from_dict(typography["body"])
        if "caption" in typography:
            ds.caption_style = TypographyStyle.from_dict(typography["caption"])
        if "metric" in typography:
            ds.metric_style = TypographyStyle.from_dict(typography["metric"])
        
        table = d.get("table", {})
        ds.table_header_bg = table.get("header_bg", "primary")
        ds.table_header_text = table.get("header_text", "#FFFFFF")
        ds.table_alternate_rows = table.get("alternate_rows", True)
        
        spacing = d.get("spacing", {})
        ds.margin_left = spacing.get("margin_left", 0.5)
        ds.margin_right = spacing.get("margin_right", 0.5)
        ds.margin_top = spacing.get("margin_top", 0.5)
        ds.margin_bottom = spacing.get("margin_bottom", 0.5)
        ds.element_gap = spacing.get("element_gap", 0.2)
        
        analysis = d.get("analysis", {})
        ds.files_analyzed = analysis.get("files_analyzed", 0)
        ds.slides_analyzed = analysis.get("slides_analyzed", 0)
        
        return ds
    
    def save(self, path: str):
        """Save to JSON file."""
        with open(path, 'w') as f:
            json.dump(self.to_dict(), f, indent=2)
    
    @classmethod
    def load(cls, path: str) -> "DesignSystem":
        """Load from JSON file."""
        with open(path, 'r') as f:
            return cls.from_dict(json.load(f))


# =============================================================================
# DESIGN SYSTEM ANALYZER
# =============================================================================

class DesignSystemAnalyzer:
    """
    Analyzes multiple PPTX files to extract a unified design system.
    
    Usage:
        analyzer = DesignSystemAnalyzer()
        design_system = analyzer.analyze(["file1.pptx", "file2.pptx", ...])
        design_system.save("my_design_system.json")
    """
    
    def __init__(self, verbose: bool = False):
        self.verbose = verbose
        self.colors: List[Tuple[str, str]] = []  # (hex, source)
        self.fonts: List[Tuple[str, float, bool, str]] = []  # (name, size, bold, context)
        self.files_analyzed = 0
        self.slides_analyzed = 0
    
    def analyze(self, pptx_paths: List[str], name: str = "Organization") -> DesignSystem:
        """
        Analyze multiple presentations and create a design system.
        
        Args:
            pptx_paths: List of paths to .pptx files
            name: Name for the design system
        
        Returns:
            DesignSystem with extracted styles
        """
        self._log(f"Analyzing {len(pptx_paths)} presentations...")
        
        # Reset
        self.colors = []
        self.fonts = []
        self.files_analyzed = 0
        self.slides_analyzed = 0
        
        for path in pptx_paths:
            if os.path.exists(path):
                try:
                    self._analyze_file(path)
                    self.files_analyzed += 1
                except Exception as e:
                    self._log(f"  Error with {path}: {e}")
        
        self._log(f"Analyzed {self.files_analyzed} files, {self.slides_analyzed} slides")
        
        return self._build_design_system(name)
    
    def _log(self, msg: str):
        if self.verbose:
            print(msg)
    
    def _analyze_file(self, path: str):
        """Analyze a single PPTX file."""
        self._log(f"  Processing: {os.path.basename(path)}")
        
        # Extract theme colors from XML
        self._extract_theme_colors(path)
        
        # Use python-pptx for content
        prs = Presentation(path)
        
        for slide in prs.slides:
            self.slides_analyzed += 1
            for shape in slide.shapes:
                self._extract_from_shape(shape)
    
    def _extract_theme_colors(self, path: str):
        """Extract colors from theme XML."""
        try:
            with zipfile.ZipFile(path, 'r') as zf:
                theme_files = [n for n in zf.namelist() if 'theme' in n.lower() and n.endswith('.xml')]
                for theme_file in theme_files:
                    with zf.open(theme_file) as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        
                        # Find srgbClr elements
                        for elem in root.iter():
                            if elem.tag.endswith('srgbClr'):
                                val = elem.get('val')
                                if val and len(val) == 6:
                                    self.colors.append((f"#{val}", "theme"))
        except:
            pass
    
    def _extract_from_shape(self, shape):
        """Extract colors and fonts from a shape."""
        # Fill color
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                    self.colors.append((f"#{shape.fill.fore_color.rgb}", "fill"))
        except:
            pass
        
        # Text properties
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Color
                    try:
                        if run.font.color and run.font.color.rgb:
                            self.colors.append((f"#{run.font.color.rgb}", "text"))
                    except:
                        pass
                    
                    # Font
                    try:
                        if run.font.name and run.font.size:
                            size_pt = run.font.size.pt
                            bold = run.font.bold or False
                            # Guess context from size
                            if size_pt >= 24:
                                context = "title"
                            elif size_pt >= 16:
                                context = "heading"
                            else:
                                context = "body"
                            self.fonts.append((run.font.name, size_pt, bold, context))
                    except:
                        pass
        
        # Recurse into groups
        if hasattr(shape, 'shapes'):
            for child in shape.shapes:
                self._extract_from_shape(child)
    
    def _build_design_system(self, name: str) -> DesignSystem:
        """Build design system from collected data."""
        ds = DesignSystem(name=name)
        ds.files_analyzed = self.files_analyzed
        ds.slides_analyzed = self.slides_analyzed
        
        # Analyze colors
        color_counts = Counter(c[0].upper() for c in self.colors)
        
        # Filter out common neutrals
        ignore = {'#FFFFFF', '#000000', '#333333', '#666666', '#999999', '#CCCCCC', '#F5F5F5'}
        brand_colors = [(c, cnt) for c, cnt in color_counts.items() if c not in ignore and cnt >= 2]
        brand_colors.sort(key=lambda x: x[1], reverse=True)
        
        if brand_colors:
            ds.primary_color = ColorValue(brand_colors[0][0], "primary")
        if len(brand_colors) > 1:
            ds.secondary_color = ColorValue(brand_colors[1][0], "secondary")
        if len(brand_colors) > 2:
            ds.accent_color = ColorValue(brand_colors[2][0], "accent")
        
        # Add more brand colors
        for i, (color, _) in enumerate(brand_colors[3:8]):
            ds.additional_colors[f"brand_{i+1}"] = ColorValue(color, f"brand_{i+1}")
        
        # Analyze fonts
        font_names = Counter(f[0] for f in self.fonts if f[0])
        primary_font = font_names.most_common(1)[0][0] if font_names else "Arial"
        
        # Calculate average sizes by context
        title_fonts = [f for f in self.fonts if f[3] == "title"]
        heading_fonts = [f for f in self.fonts if f[3] == "heading"]
        body_fonts = [f for f in self.fonts if f[3] == "body"]
        
        def avg_size(fonts, default):
            if not fonts:
                return default
            return round(sum(f[1] for f in fonts) / len(fonts), 1)
        
        def mostly_bold(fonts):
            if not fonts:
                return True
            return sum(1 for f in fonts if f[2]) > len(fonts) / 2
        
        ds.title_style = TypographyStyle(primary_font, avg_size(title_fonts, 28), mostly_bold(title_fonts))
        ds.h1_style = TypographyStyle(primary_font, avg_size([f for f in heading_fonts if f[1] >= 18], 22), True)
        ds.h2_style = TypographyStyle(primary_font, avg_size([f for f in heading_fonts if f[1] < 18], 16), True)
        ds.body_style = TypographyStyle(primary_font, avg_size(body_fonts, 11), False)
        ds.caption_style = TypographyStyle(primary_font, max(avg_size(body_fonts, 11) - 2, 8), False)
        ds.metric_style = TypographyStyle(primary_font, 36, True)
        
        return ds


# =============================================================================
# STYLE ENFORCER
# =============================================================================

class StyleEnforcer:
    """
    Enforces consistent styling from a design system.
    
    Usage:
        enforcer = StyleEnforcer(design_system)
        
        # Apply to text frame
        enforcer.apply_typography(text_frame, "title")
        
        # Get color
        rgb = enforcer.get_color_rgb("primary")
    """
    
    def __init__(self, design_system: DesignSystem):
        self.ds = design_system
        self._usage = {"colors": Counter(), "typography": Counter()}
    
    def get_color_rgb(self, name: str) -> RGBColor:
        """Get an RGBColor by name."""
        self._usage["colors"][name] += 1
        return self.ds.get_color(name).to_rgb()
    
    def get_color_hex(self, name: str) -> str:
        """Get hex color string (without #)."""
        return self.ds.get_color(name).hex.lstrip('#')
    
    def apply_typography(self, text_frame, style_name: str):
        """
        Apply a typography style to a text frame.
        
        Args:
            text_frame: python-pptx TextFrame object
            style_name: "title", "h1", "h2", "body", "caption", or "metric"
        """
        self._usage["typography"][style_name] += 1
        style = self.ds.get_typography(style_name)
        
        for para in text_frame.paragraphs:
            for run in para.runs:
                run.font.name = style.font_name
                run.font.size = Pt(style.font_size)
                run.font.bold = style.bold
                run.font.italic = style.italic
    
    def set_text_with_style(self, text_frame, text: str, style_name: str, alignment: str = "left"):
        """
        Set text and apply styling in one call.
        
        Args:
            text_frame: python-pptx TextFrame object
            text: Text to set
            style_name: Typography style name
            alignment: "left", "center", or "right"
        """
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = text
        
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
        
        self.apply_typography(text_frame, style_name)
    
    def apply_table_style(self, table):
        """Apply design system styling to a table."""
        header_bg = self.get_color_rgb(self.ds.table_header_bg)
        
        for row_idx, row in enumerate(table.rows):
            for cell in row.cells:
                if row_idx == 0:
                    # Header row
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_bg
                elif self.ds.table_alternate_rows and row_idx % 2 == 0:
                    # Alternate rows
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    
    def get_usage_report(self) -> dict:
        """Get report of what styles were used."""
        return {
            "colors": dict(self._usage["colors"]),
            "typography": dict(self._usage["typography"]),
            "design_system": self.ds.name,
        }


# =============================================================================
# STREAMLIT PAGE COMPONENT
# =============================================================================

def render_design_system_page(st):
    """
    Render a Streamlit page for managing design systems.
    
    Usage in your main app:
        from design_system_module import render_design_system_page
        
        if page == "Design System":
            render_design_system_page(st)
    """
    st.header("🎨 Design System Manager")
    
    tab1, tab2, tab3 = st.tabs(["Train", "View", "Test"])
    
    with tab1:
        st.subheader("Train from Presentations")
        st.write("Upload example presentations to learn your organization's style.")
        
        org_name = st.text_input("Organization Name", value="My Organization")
        
        uploaded_files = st.file_uploader(
            "Upload PPTX Files",
            type=["pptx"],
            accept_multiple_files=True
        )
        
        if uploaded_files and st.button("Analyze & Train", type="primary"):
            # Save uploaded files temporarily
            import tempfile
            temp_dir = tempfile.mkdtemp()
            temp_paths = []
            
            for f in uploaded_files:
                temp_path = os.path.join(temp_dir, f.name)
                with open(temp_path, 'wb') as out:
                    out.write(f.getvalue())
                temp_paths.append(temp_path)
            
            # Analyze
            with st.spinner(f"Analyzing {len(temp_paths)} files..."):
                analyzer = DesignSystemAnalyzer(verbose=False)
                ds = analyzer.analyze(temp_paths, org_name)
                st.session_state['temp_design_system'] = ds
            
            # Cleanup
            for p in temp_paths:
                os.unlink(p)
                
        # Display results if we have a temp design system
        if 'temp_design_system' in st.session_state:
            ds = st.session_state['temp_design_system']
            
            st.success(f"✅ Analyzed {ds.files_analyzed} files, {ds.slides_analyzed} slides")
            
            # Show results
            col1, col2, col3 = st.columns(3)
            with col1:
                st.color_picker("Primary", ds.primary_color.hex, disabled=True)
            with col2:
                st.color_picker("Secondary", ds.secondary_color.hex, disabled=True)
            with col3:
                if ds.accent_color:
                    st.color_picker("Accent", ds.accent_color.hex, disabled=True)
            
            st.write(f"**Font:** {ds.title_style.font_name}")
            st.write(f"**Title Size:** {ds.title_style.font_size}pt")
            st.write(f"**Body Size:** {ds.body_style.font_size}pt")
            
            # Save button
            if st.button("💾 Save Design System"):
                save_path = f"design_system_{org_name.lower().replace(' ', '_')}.json"
                ds.save(save_path)
                st.session_state['design_system'] = ds
                st.success(f"Saved to {save_path}")
                # Force reload to update View tab
                st.rerun()
    
    with tab2:
        st.subheader("Current Design System")
        
        if 'design_system' in st.session_state:
            ds = st.session_state['design_system']
            st.json(ds.to_dict())
        else:
            # Try to load from file
            json_files = list(Path(".").glob("design_system_*.json"))
            if json_files:
                selected = st.selectbox("Load Design System", json_files)
                if st.button("Load"):
                    ds = DesignSystem.load(str(selected))
                    st.session_state['design_system'] = ds
                    st.success(f"Loaded {ds.name}")
                    st.rerun()
            else:
                st.info("No design system loaded. Train one in the 'Train' tab.")
    
    with tab3:
        st.subheader("Test Generation")
        
        if 'design_system' not in st.session_state:
            st.warning("Load a design system first")
            return
        
        ds = st.session_state['design_system']
        
        test_title = st.text_input("Test Title", "Sample Site Profile")
        test_content = st.text_area("Test Content", "Point 1\nPoint 2\nPoint 3")
        
        if st.button("Generate Test Slide"):
            from pptx import Presentation as PptxPresentation
            
            prs = PptxPresentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            
            enforcer = StyleEnforcer(ds)
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(ds.margin_left),
                Inches(ds.margin_top),
                Inches(12),
                Inches(1)
            )
            enforcer.set_text_with_style(title_box.text_frame, test_title, "title")
            
            # Content
            content_box = slide.shapes.add_textbox(
                Inches(ds.margin_left),
                Inches(1.5),
                Inches(12),
                Inches(5)
            )
            tf = content_box.text_frame
            for i, line in enumerate(test_content.split("\n")):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
            enforcer.apply_typography(tf, "body")
            
            # Save
            output_path = "test_output.pptx"
            prs.save(output_path)
            
            with open(output_path, 'rb') as f:
                st.download_button(
                    "📥 Download Test PPTX",
                    f.read(),
                    file_name="test_output.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            
            st.success("Generated! Download to preview styling.")


# =============================================================================
# QUICK TEST
# =============================================================================

if __name__ == "__main__":
    print("Design System Module")
    print("=" * 50)
    
    # Create a sample design system
    ds = DesignSystem(name="Test Organization")
    ds.primary_color = ColorValue("#E31937")  # Red
    ds.secondary_color = ColorValue("#1A2B4A")  # Navy
    
    print(f"Created design system: {ds.name}")
    print(f"Primary color: {ds.primary_color.hex}")
    print(f"Title font: {ds.title_style.font_name} {ds.title_style.font_size}pt")
    
    # Save test
    ds.save("test_design_system.json")
    print("Saved to test_design_system.json")
    
    # Load test
    loaded = DesignSystem.load("test_design_system.json")
    print(f"Loaded: {loaded.name}")
    
    # Cleanup
    os.unlink("test_design_system.json")
    print("Test complete!")
