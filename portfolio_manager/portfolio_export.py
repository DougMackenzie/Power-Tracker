"""
Portfolio Export Module
=======================
Generates a comprehensive PowerPoint deck for a selected portfolio of sites.
Includes:
1. Portfolio Summary (Metrics, Charts, Rankings)
2. State Analysis (if applicable)
3. Individual Site Profiles (Site Profile, Capacity, Critical Path, etc.)
"""

import copy
import io
from typing import List, Dict, Any
from datetime import datetime
import matplotlib.pyplot as plt
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = None

from .pptx_export import (
    JLL_COLORS, SiteProfileData, CapacityTrajectory, PhaseData, ScoreAnalysis,
    TEMPLATE_SLIDES, ExportConfig,
    generate_capacity_trajectory_chart, generate_critical_path_chart,
    generate_score_radar_chart, generate_score_summary_chart,
    generate_market_analysis_chart,
    populate_site_profile_table, update_overview_textbox,
    find_and_replace_text, build_replacements,
    add_header_bar, add_footer, set_slide_background_white,
    add_critical_path_text, add_score_breakdown_text, add_market_text
)
from .program_tracker import calculate_portfolio_summary, ProgramTrackerData


def copy_slide_from_external(source_slide, dest_prs):
    """
    Copy a slide from a source presentation to the destination presentation.
    Handles TextBoxes, AutoShapes, Pictures, and Tables.
    """
    # Create a blank slide
    blank_layout = dest_prs.slide_layouts[6] 
    dest_slide = dest_prs.slides.add_slide(blank_layout)
    
    # Copy shapes
    for shape in source_slide.shapes:
        # 1. Pictures
        if shape.shape_type == 13: # PICTURE
            try:
                blob = shape.image.blob
                dest_slide.shapes.add_picture(io.BytesIO(blob), shape.left, shape.top, shape.width, shape.height)
            except Exception as e:
                print(f"Failed to copy picture: {e}")
        
        # 2. Tables
        elif shape.shape_type == 19: # TABLE
            try:
                rows = len(shape.table.rows)
                cols = len(shape.table.columns)
                new_table = dest_slide.shapes.add_table(
                    rows, cols, shape.left, shape.top, shape.width, shape.height
                ).table
                
                for r in range(rows):
                    for c in range(cols):
                        source_cell = shape.table.cell(r, c)
                        dest_cell = new_table.cell(r, c)
                        
                        # Copy text and simple formatting
                        dest_cell.text = source_cell.text
                        
                        # Copy paragraph alignment/font from first paragraph if exists
                        if source_cell.text_frame.paragraphs:
                            p_source = source_cell.text_frame.paragraphs[0]
                            p_dest = dest_cell.text_frame.paragraphs[0]
                            p_dest.alignment = p_source.alignment
                            
                            if p_source.runs:
                                r_source = p_source.runs[0]
                                if p_dest.runs:
                                    r_dest = p_dest.runs[0]
                                    r_dest.font.size = r_source.font.size
                                    r_dest.font.bold = r_source.font.bold
                                    try:
                                        r_dest.font.color.rgb = r_source.font.color.rgb
                                    except:
                                        pass
            except Exception as e:
                print(f"Failed to copy table: {e}")

        # 3. TextBoxes and AutoShapes
        elif shape.shape_type == MSO_SHAPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE.TEXT_BOX:
            try:
                new_shape = dest_slide.shapes.add_shape(
                    shape.auto_shape_type,
                    shape.left, shape.top, shape.width, shape.height
                )
                
                # Copy fill (basic)
                if shape.fill.type == 1: # Solid
                    try:
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                    except:
                        pass
                
                # Copy text
                if shape.has_text_frame:
                    new_shape.text_frame.clear()
                    for paragraph in shape.text_frame.paragraphs:
                        new_p = new_shape.text_frame.add_paragraph()
                        new_p.text = paragraph.text
                        new_p.alignment = paragraph.alignment
                        
                        # Copy runs
                        for run in paragraph.runs:
                            new_run = new_p.add_run()
                            new_run.text = run.text
                            new_run.font.name = run.font.name
                            new_run.font.size = run.font.size
                            new_run.font.bold = run.font.bold
                            new_run.font.italic = run.font.italic
                            try:
                                new_run.font.color.rgb = run.font.color.rgb
                            except:
                                pass
            except Exception as e:
                print(f"Failed to copy shape: {e}")




def prepare_site_for_export(site_data: Dict) -> Dict:
    """
    Parse JSON string fields in site_data into Python objects.
    This ensures charts and detailed slides can access structured data.
    """
    import json
    
    # Fields that might be JSON strings
    json_fields = [
        'phases_json', 'schedule_json', 'capacity_trajectory',
        'infrastructure', 'score_analysis', 'market_analysis'
    ]
    
    for field in json_fields:
        if field in site_data and isinstance(site_data[field], str):
            try:
                site_data[field] = json.loads(site_data[field])
            except:
                pass  # Not valid JSON, leave as-is
    
    return site_data


def get_profile_data(site_data: Dict):
    """
    Robustly extract SiteProfileData from various formats.
    Handles: existing object, dict, JSON string, or maps raw site_data.
    """
    from .pptx_export import SiteProfileData, map_app_to_profile
    import json
    
    # 1. Already a SiteProfileData object
    if 'profile' in site_data:
        p = site_data['profile']
        if hasattr(p, 'overview') and hasattr(p, 'to_description_dict'):
            return p
        # 2. Dictionary format
        elif isinstance(p, dict):
            return SiteProfileData.from_dict(p)
        # 3. JSON string
        elif isinstance(p, str):
            try:
                p_dict = json.loads(p)
                return SiteProfileData.from_dict(p_dict)
            except:
                pass
    
    # 4. Fallback: map from raw site_data
    try:
        return map_app_to_profile(site_data)
    except:
        return None


def generate_portfolio_export(
    sites: Dict[str, Dict],
    template_path: str,
    output_path: str,
    config: ExportConfig = None
) -> str:
    """
    Generate portfolio export by creating individual presentations for each site,
    then merging them together. This guarantees exact fidelity to single-site exports.
    """
    import tempfile
    import os
    from .pptx_export import export_site_to_pptx
    
    if not Presentation:
        raise ImportError("python-pptx is required")
    
    if config is None:
        config = ExportConfig()
    
    print(f"[DEBUG] Generating Portfolio Export for {len(sites)} sites using merge strategy")
    
    # 1. Create master presentation (start with template)
    master_prs = Presentation(template_path)
    
    # 2. Update Title Slide
    title_slide = master_prs.slides[0]
    if title_slide.shapes.title:
        title_slide.shapes.title.text = "Portfolio Overview"
    
    subtitle_text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    if len(title_slide.placeholders) > 1:
        try:
            title_slide.placeholders[1].text = subtitle_text
        except:
            pass
    
    # 3. Add Portfolio Summary Slides
    add_portfolio_metrics_slide(master_prs, sites, index=1)
    add_portfolio_ranking_slide(master_prs, sites, index=2)
    
    # 4. Generate individual presentations for each site and merge
    temp_files = []
    
    for site_id, site_data in sites.items():
        site_name = site_data.get('name', site_id)
        print(f"[DEBUG] Generating individual export for: {site_name}")
        
        try:
            # Generate individual presentation using the proven export function
            temp_path = tempfile.mktemp(suffix='.pptx')
            temp_files.append(temp_path)
            
            export_site_to_pptx(site_data, template_path, temp_path, config)
            
            # Load the individual presentation
            site_prs = Presentation(temp_path)
            
            # Copy all slides EXCEPT the title slide (index 0) from the individual presentation
            # Individual exports have: Title, Profile, Boundary, Topo, Charts...
            # We want to copy everything except the title
            for slide_idx in range(1, len(site_prs.slides)):
                source_slide = site_prs.slides[slide_idx]
                
                # Get the layout from source
                slide_layout = source_slide.slide_layout
                
                # Add a new slide with the same layout
                new_slide = master_prs.slides.add_slide(slide_layout)
                
                # Copy all shapes from source to new slide
                copy_slide_content(source_slide, new_slide)
            
            print(f"[DEBUG] Merged {len(site_prs.slides) - 1} slides from {site_name}")
            
        except Exception as e:
            print(f"[ERROR] Failed to export/merge {site_name}: {e}")
            import traceback
            traceback.print_exc()
    
    # 5. Add Thank You slide at the end
    try:
        blank_layout = master_prs.slide_layouts[6]  # Blank layout
        thank_you_slide = master_prs.slides.add_slide(blank_layout)
        
        # Add centered text
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1)
        
        textbox = thank_you_slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Thank You"
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = 1  # Center
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
    except:
        pass
    
    # 6. Clean up temp files
    for temp_file in temp_files:
        try:
            if os.path.exists(temp_file):
                os.unlink(temp_file)
        except:
            pass
    
    # 7. Save master portfolio
    master_prs.save(output_path)
    print(f"[DEBUG] Portfolio export saved to: {output_path}")
    
    return output_path


def copy_slide_content(source_slide, dest_slide):
    """
    Copy all shapes from source slide to destination slide.
    Uses deep copy for complete fidelity.
    """
    import copy
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    # Clear any existing shapes in destination
    for shape in list(dest_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)
    
    # Copy all shapes from source
    for shape in source_slide.shapes:
        try:
            # Use deep copy of the XML element for complete fidelity
            new_el = copy.deepcopy(shape.element)
            dest_slide.shapes._spTree.append(new_el)
        except Exception as e:
            print(f"[WARNING] Failed to copy shape: {e}")
    
    return dest_slide


def duplicate_slide_in_place(prs, index):
    """
    Duplicate a slide within the same presentation.
    Uses the SAME LAYOUT as the source slide to preserve structure.
    """
    source = prs.slides[index]
    layout = source.slide_layout
    dest = prs.slides.add_slide(layout)
    
    # CRITICAL: Clear existing shapes from the new slide!
    for shape in list(dest.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Copy shapes
    for shape in source.shapes:
        copy_shape(shape, dest.shapes)
            
    return dest


def copy_shape(shape, shapes_collection, group_offset=(0,0)):
    """
    Recursively copy a shape to a shapes collection.
    group_offset: (left, top) adjustment for shapes inside groups.
    """
    # 1. Groups (Type 6)
    if shape.shape_type == 6: 
        # Create a new group
        # python-pptx doesn't support creating groups easily from scratch via API
        # We have to copy members individually.
        # Note: This "ungroups" the group in the copy, but preserves visual appearance.
        # To truly preserve grouping would require XML manipulation which is risky.
        # We will copy members as individual shapes.
        for member in shape.shapes:
            # Adjust position by group's position
            # Actually, member.left/top are absolute in python-pptx
            copy_shape(member, shapes_collection)
        return

    # 2. Pictures (Type 13)
    if shape.shape_type == 13: 
        try:
            blob = shape.image.blob
            shapes_collection.add_picture(io.BytesIO(blob), shape.left, shape.top, shape.width, shape.height)
        except: pass
        return

    # 3. Tables (Type 19)
    if shape.shape_type == 19:
        try:
            rows = len(shape.table.rows)
            cols = len(shape.table.columns)
            new_table_shape = shapes_collection.add_table(rows, cols, shape.left, shape.top, shape.width, shape.height)
            new_table = new_table_shape.table
            
            # Copy Column Widths
            for i in range(cols):
                new_table.columns[i].width = shape.table.columns[i].width
            # Copy Row Heights
            for i in range(rows):
                new_table.rows[i].height = shape.table.rows[i].height
                
            # Copy content
            for r in range(rows):
                for c in range(cols):
                    source_cell = shape.table.cell(r, c)
                    dest_cell = new_table.cell(r, c)
                    copy_text_frame(source_cell.text_frame, dest_cell.text_frame)
                    
                    # Copy cell fill
                    try:
                        if source_cell.fill.type == 1: # Solid
                            dest_cell.fill.solid()
                            dest_cell.fill.fore_color.rgb = source_cell.fill.fore_color.rgb
                    except: pass
                    
                    # Copy borders (simplified - copying all borders is complex)
        except: pass
        return

    # 4. AutoShapes (1) / TextBoxes (17) / Placeholders (14)
    # Treat placeholders as shapes to preserve their content
    if shape.shape_type in [1, 17, 14]:
        try:
            shape_type = shape.auto_shape_type if hasattr(shape, 'auto_shape_type') else 1 
            new_shape = shapes_collection.add_shape(shape_type, shape.left, shape.top, shape.width, shape.height)
            
            # Copy Text
            if shape.has_text_frame:
                copy_text_frame(shape.text_frame, new_shape.text_frame)
            
            # Copy Fill
            if shape.fill.type == 1:
                new_shape.fill.solid()
                try:
                    new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
                except: pass
                
            # Copy Line
            if shape.line.fill.type == 1:
                try:
                    new_shape.line.color.rgb = shape.line.color.rgb
                    new_shape.line.width = shape.line.width
                except: pass
        except: pass
        return


def copy_text_frame(source_tf, dest_tf):
    """Copy text frame content and formatting."""
    dest_tf.clear()
    
    # Copy vertical alignment
    if source_tf.vertical_anchor:
        dest_tf.vertical_anchor = source_tf.vertical_anchor
    
    for p in source_tf.paragraphs:
        new_p = dest_tf.add_paragraph()
        new_p.text = p.text
        new_p.alignment = p.alignment
        new_p.level = p.level
        
        # Copy paragraph spacing
        new_p.space_before = p.space_before
        new_p.space_after = p.space_after
        new_p.line_spacing = p.line_spacing
        
        # Copy runs
        for r in p.runs:
            new_r = new_p.add_run()
            new_r.text = r.text
            new_r.font.size = r.font.size
            new_r.font.bold = r.font.bold
            new_r.font.italic = r.font.italic
            new_r.font.underline = r.font.underline
            new_r.font.name = r.font.name # Copy font name!
            
            try:
                new_r.font.color.rgb = r.font.color.rgb
            except: pass
            try:
                if r.font.color.type == 2: # Theme color
                    new_r.font.color.theme_color = r.font.color.theme_color
            except: pass


def replace_images_with_placeholders(slide, site_data, label="Map Placeholder"):
    """Replace images with gray placeholders."""
    
    shapes_to_replace = []
    for shape in slide.shapes:
        # Identify images to replace
        # Type 13 = Picture
        if shape.shape_type == 13:
            # Heuristic: 
            # 1. If label is "Map Placeholder" (Profile Slide), replace images on right (>4 inches)
            # 2. If label is "Site Boundary" or "Topography", replace ANY large image (>4 inches width)
            
            if label == "Map Placeholder":
                if shape.left > Inches(4):
                    shapes_to_replace.append(shape)
            else:
                # For boundary/topo, replace main images
                if shape.width > Inches(4):
                    shapes_to_replace.append(shape)
                
    for shape in shapes_to_replace:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        # Remove
        sp = shape._element
        sp.getparent().remove(sp)
        
        # Add Placeholder
        rect = slide.shapes.add_shape(1, left, top, width, height) # 1 = Rectangle
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
        rect.line.color.rgb = RGBColor(200, 200, 200)
        
        tf = rect.text_frame
        tf.text = label
        tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_portfolio_metrics_slide(prs, sites, index=None):
    """Add slide with portfolio metrics and charts."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Move slide if index is specified
    if index is not None:
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(index, slides[-1])
    
    # Header
    add_header_bar(slide, "Portfolio Summary", Inches, Pt, RGBColor)
    
    # Calculate metrics
    site_list = []
    for s in sites.values():
        site_list.append(s)
    summary = calculate_portfolio_summary(site_list)
    
    # Metrics Text
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(3.0)
    height = Inches(5.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    metrics = [
        ("Total Sites", f"{summary['site_count']}"),
        ("Total Fee Potential", f"${summary['total_potential']:,.0f}"),
        ("Weighted Pipeline", f"${summary['total_weighted']:,.0f}"),
        ("Avg Probability", f"{summary['avg_probability']:.1%}"),
    ]
    
    for label, value in metrics:
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 43, 74)
        
        p = tf.add_paragraph()
        p.text = value
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(46, 125, 50) # Green
        p.space_after = Pt(20)

    # Charts (Generated via Matplotlib)
    # 1. Pipeline by Stage
    fig, ax = plt.subplots(figsize=(6, 4))
    stages = list(summary['by_stage'].keys())
    values = [sum(s['weighted'] for s in summary['by_stage'][stage]) for stage in stages]
    
    ax.bar(stages, values, color=[JLL_COLORS['teal'], JLL_COLORS['amber'], JLL_COLORS['light_blue'], JLL_COLORS['green']])
    ax.set_title('Weighted Value by Stage', fontsize=12, fontweight='bold', color=JLL_COLORS['dark_blue'])
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x/1e6:.0f}M'))
    plt.xticks(rotation=45)
    plt.tight_layout()
    
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=150)
    img_stream.seek(0)
    plt.close()
    
    slide.shapes.add_picture(img_stream, Inches(4.0), Inches(1.5), Inches(5.0), Inches(3.5))
    

def add_portfolio_ranking_slide(prs, sites, index=None):
    """Add slide with ranking table."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Move slide if index is specified
    if index is not None:
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(index, slides[-1])
    
    add_header_bar(slide, "Portfolio Rankings", Inches, Pt, RGBColor)
    
    # Sort sites by score
    ranked_sites = []
    for sid, s in sites.items():
        score = 0
        if 'profile' in s and isinstance(s['profile'], SiteProfileData):
            score = s['profile'].ratings.get('overall_score', 0) # This might be wrong place
        # Try to get score from scoring engine output if available
        # For now, let's assume 'score' key exists or calculate it
        # Simplified:
        ranked_sites.append((s.get('name', sid), s.get('total_fee_potential', 0), s.get('probability', 0)))
        
    # Sort by Fee * Prob (Weighted)
    ranked_sites.sort(key=lambda x: x[1]*x[2], reverse=True)
    
    # Table
    rows = min(len(ranked_sites) + 1, 11) # Header + top 10
    cols = 4
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(11.3)
    height = Inches(0.5 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header
    headers = ["Site Name", "Fee Potential", "Probability", "Weighted Value"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(26, 43, 74)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        
    # Data
    for i, (name, fee, prob) in enumerate(ranked_sites[:10]):
        row = i + 1
        table.cell(row, 0).text = name
        table.cell(row, 1).text = f"${fee:,.0f}"
        table.cell(row, 2).text = f"{prob:.1%}"
        table.cell(row, 3).text = f"${fee*prob:,.0f}"


def add_capacity_slide(prs, site_data, replacements):
    """Add Capacity Trajectory slide."""
    from .pptx_export import generate_capacity_trajectory_chart, CapacityTrajectory, PhaseData, convert_phase_data
    import tempfile
    import os
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, "Capacity Trajectory", Inches, Pt, RGBColor)
    
    # Generate Chart
    traj_data = site_data.get('capacity_trajectory', site_data.get('schedule', {}))
    trajectory = (CapacityTrajectory.from_dict(traj_data) if traj_data else
                  CapacityTrajectory.generate_default(
                      site_data.get('target_mw', 600),
                      site_data.get('phase1_mw'),
                      site_data.get('start_year', 2028)))
    
    # Extract Phases
    phases = []
    phase_data = site_data.get('phases', [])
    if phase_data:
        for idx, pd in enumerate(phase_data, 1):
            if isinstance(pd, dict):
                converted_pd = convert_phase_data(pd)
                if 'phase_num' not in converted_pd:
                    converted_pd['phase_num'] = idx
                phases.append(PhaseData(**converted_pd))
            else:
                phases.append(pd)
                      
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        chart_path = tmp.name
        
    try:
        generate_capacity_trajectory_chart(
            trajectory, 
            site_data.get('name', 'Site'), 
            chart_path, 
            phases=phases
        )
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
    except Exception as e:
        print(f"[ERROR] Failed to generate capacity chart: {e}")
    finally:
        if os.path.exists(chart_path):
            os.unlink(chart_path)


def add_infrastructure_slide(prs, site_data, replacements):
    """Add Infrastructure slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, "Infrastructure Analysis", Inches, Pt, RGBColor)
    
    # Placeholder content
    left = Inches(1)
    top = Inches(2)
    width = Inches(11)
    height = Inches(4)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Infrastructure Details"
    p = tf.add_paragraph()
    p.text = f"Substation: {site_data.get('substation', 'TBD')}"
    p = tf.add_paragraph()
    p.text = f"Transmission Line: {site_data.get('transmission_line', 'TBD')}"


def add_market_slide(prs, site_data, replacements):
    """Add Market Analysis slide."""
    from .pptx_export import generate_market_analysis_chart, add_market_text
    import tempfile
    import os
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, "Market Analysis", Inches, Pt, RGBColor)
    
    # Chart
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        chart_path = tmp.name
        
    try:
        generate_market_analysis_chart(
            site_data, 
            site_data.get('name', 'Site'), 
            chart_path
        )
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.5))
    except Exception as e:
        print(f"[ERROR] Failed to generate market chart: {e}")
    finally:
        if os.path.exists(chart_path):
            os.unlink(chart_path)
        
    # Text
    add_market_text(slide, site_data)


def add_score_slide(prs, site_data, replacements):
    """Add Score Analysis slide."""
    from .pptx_export import generate_score_summary_chart, add_score_breakdown_text, ScoreAnalysis
    import tempfile
    import os
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header_bar(slide, "Site Scoring Analysis", Inches, Pt, RGBColor)
    
    # Score Data
    score_data = None
    if 'profile' in site_data and hasattr(site_data['profile'], 'ratings'):
         score_data = ScoreAnalysis.from_dict(site_data['profile'].ratings)
    
    if score_data:
        # Chart
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            chart_path = tmp.name
            
        try:
            generate_score_summary_chart(
                score_data, 
                site_data.get('name', 'Site'), 
                chart_path
            )
            slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.5))
        except Exception as e:
            print(f"[ERROR] Failed to generate score chart: {e}")
        finally:
            if os.path.exists(chart_path):
                os.unlink(chart_path)
            
        # Text
        add_score_breakdown_text(slide, score_data)
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(1))
        txBox.text_frame.text = "No scoring data available for this site."
