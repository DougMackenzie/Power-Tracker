"""
Generate Sample Presentations for Testing
==========================================
Creates 3 sample PPTX files with consistent branding
to test the design system analyzer.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors (simulating a consistent corporate style)
NAVY = RGBColor(0x1A, 0x2B, 0x4A)
RED = RGBColor(0xE3, 0x19, 0x37)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "Arial"


def create_title_slide(prs, title, subtitle=""):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_NAME
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE


def create_content_slide(prs, title, bullets):
    """Create a content slide with bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(12)


def create_metrics_slide(prs, title, metrics):
    """Create a slide with key metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Metrics
    num_metrics = len(metrics)
    metric_width = (12.3 - 0.5 * (num_metrics - 1)) / num_metrics
    
    for i, metric in enumerate(metrics):
        left = 0.5 + i * (metric_width + 0.5)
        
        # Value
        value_box = slide.shapes.add_textbox(Inches(left), Inches(2.5), Inches(metric_width), Inches(1.2))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['value']
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_NAME
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RED
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(left), Inches(3.8), Inches(metric_width), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['label']
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY


def create_table_slide(prs, title, headers, data):
    """Create a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    rows = len(data) + 1
    cols = len(headers)
    
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.3), Inches(rows * 0.5)).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            para.font.name = FONT_NAME
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = WHITE
    
    # Data rows
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for para in cell.text_frame.paragraphs:
                para.font.name = FONT_NAME
                para.font.size = Pt(11)
                para.font.color.rgb = GRAY


def generate_sample_1():
    """Generate Site Profile presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    create_title_slide(prs, "Rogers County Site Profile", "Oklahoma - 500 MW Opportunity")
    
    create_content_slide(prs, "Site Overview", [
        "Location: Rogers County, Oklahoma",
        "Target Capacity: 500 MW",
        "Land Status: PSA in negotiation",
        "Interconnection: Queue position #247 at Tulsa North 345kV",
        "Water: Municipal water available within 2 miles",
        "Fiber: Multiple carriers present along Highway 66"
    ])
    
    create_metrics_slide(prs, "Key Metrics", [
        {"value": "500", "label": "Target MW"},
        {"value": "2028", "label": "Energization"},
        {"value": "$45M", "label": "Land Cost"},
        {"value": "320", "label": "Acres"}
    ])
    
    create_table_slide(prs, "Timeline Milestones", 
        ["Milestone", "Target Date", "Status"],
        [
            ["Site Control", "Q2 2025", "In Progress"],
            ["SIS Complete", "Q3 2025", "Pending"],
            ["Facilities Study", "Q1 2026", "Scheduled"],
            ["IA Execution", "Q3 2026", "Future"],
            ["Energization", "Q2 2028", "Future"]
        ]
    )
    
    create_content_slide(prs, "Risk Factors", [
        "Queue congestion - PSO processing delays reported",
        "Transformer lead time extending to 4+ years",
        "County zoning requires special use permit",
        "Competing projects in area may impact capacity allocation"
    ])
    
    prs.save("sample_site_profile.pptx")
    print("Created: sample_site_profile.pptx")


def generate_sample_2():
    """Generate Portfolio Summary presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    create_title_slide(prs, "Portfolio Summary", "Q4 2025 Update")
    
    create_metrics_slide(prs, "Portfolio Overview", [
        {"value": "12", "label": "Active Sites"},
        {"value": "4.2 GW", "label": "Total Capacity"},
        {"value": "6", "label": "States"},
        {"value": "$1.2B", "label": "Total Investment"}
    ])
    
    create_table_slide(prs, "Site Status Summary",
        ["Site", "State", "MW", "Stage", "Target Date"],
        [
            ["Rogers County", "OK", "500", "Development", "Q2 2028"],
            ["Mayes Industrial", "OK", "750", "Due Diligence", "Q4 2028"],
            ["Pinal Solar", "AZ", "400", "Construction", "Q1 2026"],
            ["Navajo Gateway", "NM", "600", "Permitting", "Q3 2027"],
            ["West Texas Hub", "TX", "1,000", "Development", "Q1 2029"]
        ]
    )
    
    create_content_slide(prs, "Q4 Priorities", [
        "Close Rogers County PSA by December 15",
        "Complete Mayes Industrial site assessment",
        "Submit Pinal Solar construction permits",
        "Finalize transformer procurement for 3 sites",
        "Begin Navajo Gateway environmental review"
    ])
    
    create_content_slide(prs, "Key Risks & Mitigations", [
        "Supply Chain: Transformer lead times extending - secured slots with 2 manufacturers",
        "Regulatory: ERCOT interconnection delays - pursuing fast-track options",
        "Market: Power prices volatile - hedging strategies in place",
        "Execution: Resource constraints - adding 2 project managers"
    ])
    
    prs.save("sample_portfolio_summary.pptx")
    print("Created: sample_portfolio_summary.pptx")


def generate_sample_3():
    """Generate Market Analysis presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    create_title_slide(prs, "Oklahoma Power Market", "Data Center Development Opportunity")
    
    create_content_slide(prs, "Market Drivers", [
        "Hyperscaler demand exceeding 2 GW annually in Southwest region",
        "Oklahoma positioned as low-cost alternative to Texas ERCOT",
        "SPP capacity surplus provides competitive pricing through 2027",
        "State incentives: 5-year property tax abatement for data centers",
        "Renewable energy: Wind capacity supports sustainability mandates"
    ])
    
    create_metrics_slide(prs, "Oklahoma Advantages", [
        {"value": "$0.045", "label": "Avg $/kWh"},
        {"value": "15%", "label": "Below TX"},
        {"value": "2.8 GW", "label": "Wind Capacity"},
        {"value": "AA", "label": "Grid Rating"}
    ])
    
    create_table_slide(prs, "Utility Comparison",
        ["Utility", "Territory", "Capacity", "Lead Time", "Rating"],
        [
            ["PSO", "East OK", "High", "24-36 mo", "A"],
            ["OG&E", "Central OK", "Medium", "18-24 mo", "A+"],
            ["SWEPCO", "SE OK", "Low", "30-42 mo", "B+"]
        ]
    )
    
    create_content_slide(prs, "Recommended Strategy", [
        "Focus on PSO territory for near-term opportunities",
        "Target 345kV substations with headroom > 300 MW",
        "Pursue sites within 10 miles of existing transmission",
        "Engage OG&E for 2027+ pipeline development",
        "Avoid SWEPCO territory due to extended timelines"
    ])
    
    prs.save("sample_market_analysis.pptx")
    print("Created: sample_market_analysis.pptx")


if __name__ == "__main__":
    print("Generating sample presentations...")
    print("=" * 50)
    generate_sample_1()
    generate_sample_2()
    generate_sample_3()
    print("=" * 50)
    print("Done! Created 3 sample PPTX files for testing.")
